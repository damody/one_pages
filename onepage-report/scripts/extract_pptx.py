#!/usr/bin/env python3
"""
從 PPTX 檔案抽取文字與圖片

用法：
    python extract_pptx.py input.pptx output_dir/ [--slides "1-3,7,10-12"]

輸出：
    output_dir/text.md      - 投影片文字（含 slide# 與 shape# 標注）
    output_dir/images/      - 抽出的圖片
    output_dir/summary.md   - 投影片摘要（每頁標題）

範例：
    python extract_pptx.py presentation.pptx ./extracted/
    python extract_pptx.py presentation.pptx ./extracted/ --slides "1-5,8,10-15"
"""

import argparse
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def parse_slide_range(range_str: str, total_slides: int) -> list[int]:
    """
    解析投影片範圍字串

    範例：
        "1-3,7,10-12" → [1, 2, 3, 7, 10, 11, 12]
        "1-5" → [1, 2, 3, 4, 5]
        "3" → [3]

    Args:
        range_str: 範圍字串
        total_slides: 總投影片數量

    Returns:
        投影片編號列表（1-based）
    """
    if not range_str:
        return list(range(1, total_slides + 1))

    slides = set()
    parts = range_str.replace(" ", "").split(",")

    for part in parts:
        if "-" in part:
            start, end = part.split("-", 1)
            start = int(start)
            end = int(end)
            for i in range(start, min(end + 1, total_slides + 1)):
                if 1 <= i <= total_slides:
                    slides.add(i)
        else:
            num = int(part)
            if 1 <= num <= total_slides:
                slides.add(num)

    return sorted(slides)


def extract_shape_text(shape) -> str:
    """從 shape 中抽取文字"""
    text_parts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in paragraph.runs)
            if para_text.strip():
                text_parts.append(para_text.strip())

    return "\n".join(text_parts)


def get_slide_title(slide) -> str:
    """取得投影片標題"""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()

    # 嘗試從第一個文字方塊取得標題
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # 取第一行作為標題
                return text.split("\n")[0][:50]

    return "(無標題)"


def extract_pptx(pptx_path: str, output_dir: str, slide_range: str = None) -> dict:
    """
    從 PPTX 抽取內容

    Args:
        pptx_path: PPTX 檔案路徑
        output_dir: 輸出目錄
        slide_range: 投影片範圍（如 "1-3,7"）

    Returns:
        抽取結果摘要
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)

    if not pptx_path.exists():
        raise FileNotFoundError(f"找不到檔案：{pptx_path}")

    # 建立輸出目錄
    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    # 載入 PPTX
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)

    # 解析投影片範圍
    selected_slides = parse_slide_range(slide_range, total_slides)

    # 收集內容
    text_content = []
    summary_content = []
    image_count = 0

    text_content.append(f"# PPTX 內容抽取：{pptx_path.name}")
    text_content.append(f"")
    text_content.append(f"- 總投影片數：{total_slides}")
    text_content.append(f"- 抽取範圍：{selected_slides}")
    text_content.append(f"")
    text_content.append("---")
    text_content.append("")

    summary_content.append(f"# 投影片摘要：{pptx_path.name}")
    summary_content.append("")
    summary_content.append("| 頁碼 | 標題 |")
    summary_content.append("|------|------|")

    for slide_num in selected_slides:
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]
        title = get_slide_title(slide)

        # 摘要
        summary_content.append(f"| {slide_num} | {title} |")

        # 詳細內容
        text_content.append(f"## Slide {slide_num}：{title}")
        text_content.append("")

        shape_num = 0
        for shape in slide.shapes:
            shape_num += 1

            # 抽取文字
            if shape.has_text_frame:
                text = extract_shape_text(shape)
                if text:
                    text_content.append(f"### Shape {shape_num}")
                    text_content.append(f"[來源：slide={slide_num}, shape={shape_num}]")
                    text_content.append("")
                    text_content.append(text)
                    text_content.append("")

            # 抽取圖片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_ext = image.ext
                    image_filename = f"slide{slide_num}_shape{shape_num}.{image_ext}"
                    image_path = images_dir / image_filename

                    with open(image_path, "wb") as f:
                        f.write(image.blob)

                    image_count += 1
                    text_content.append(f"### 圖片：{image_filename}")
                    text_content.append(f"[來源：slide={slide_num}, shape={shape_num}]")
                    text_content.append(f"")
                    text_content.append(f"![{image_filename}](images/{image_filename})")
                    text_content.append("")
                except Exception as e:
                    text_content.append(f"### 圖片抽取失敗")
                    text_content.append(f"[來源：slide={slide_num}, shape={shape_num}]")
                    text_content.append(f"錯誤：{str(e)}")
                    text_content.append("")

        text_content.append("---")
        text_content.append("")

    # 寫入檔案
    text_path = output_dir / "text.md"
    with open(text_path, "w", encoding="utf-8") as f:
        f.write("\n".join(text_content))

    summary_path = output_dir / "summary.md"
    with open(summary_path, "w", encoding="utf-8") as f:
        f.write("\n".join(summary_content))

    result = {
        "pptx_file": str(pptx_path),
        "total_slides": total_slides,
        "extracted_slides": len(selected_slides),
        "image_count": image_count,
        "output_files": {
            "text": str(text_path),
            "summary": str(summary_path),
            "images_dir": str(images_dir)
        }
    }

    return result


def main():
    parser = argparse.ArgumentParser(
        description="從 PPTX 檔案抽取文字與圖片",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例：
    python extract_pptx.py presentation.pptx ./output/
    python extract_pptx.py presentation.pptx ./output/ --slides "1-5,8,10-15"
    python extract_pptx.py presentation.pptx ./output/ --slides "1-3"

投影片範圍格式：
    "1-3"       選取第 1, 2, 3 頁
    "1,3,5"     選取第 1, 3, 5 頁
    "1-3,7,10-12" 選取第 1, 2, 3, 7, 10, 11, 12 頁
        """
    )

    parser.add_argument("pptx_path", help="PPTX 檔案路徑")
    parser.add_argument("output_dir", help="輸出目錄")
    parser.add_argument("--slides", "-s", help="投影片範圍（如 '1-3,7,10-12'）")
    parser.add_argument("--list", "-l", action="store_true", help="只列出投影片清單，不抽取")

    args = parser.parse_args()

    # 只列出投影片清單
    if args.list:
        prs = Presentation(args.pptx_path)
        print(f"檔案：{args.pptx_path}")
        print(f"總投影片數：{len(prs.slides)}")
        print("")
        print("| 頁碼 | 標題 |")
        print("|------|------|")
        for i, slide in enumerate(prs.slides, 1):
            title = get_slide_title(slide)
            print(f"| {i} | {title} |")
        return

    # 抽取內容
    try:
        result = extract_pptx(args.pptx_path, args.output_dir, args.slides)

        print(f"抽取完成！")
        print(f"")
        print(f"來源：{result['pptx_file']}")
        print(f"總投影片：{result['total_slides']} 頁")
        print(f"已抽取：{result['extracted_slides']} 頁")
        print(f"圖片數：{result['image_count']} 張")
        print(f"")
        print(f"輸出檔案：")
        print(f"  - {result['output_files']['text']}")
        print(f"  - {result['output_files']['summary']}")
        print(f"  - {result['output_files']['images_dir']}/")

    except FileNotFoundError as e:
        print(f"錯誤：{e}")
        sys.exit(1)
    except Exception as e:
        print(f"抽取失敗：{e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
