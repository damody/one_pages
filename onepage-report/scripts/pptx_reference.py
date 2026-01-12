#!/usr/bin/env python3
"""
python-pptx API Reference for AI-generated PPTX code.
This file demonstrates common patterns - AI should adapt based on content.

MTK 風格：
- 背景色（米色）
- 圓角矩形 content box
- 顏色區分區塊（藍/橙/綠/紫）
- 表格支援
- 統一字體 (Microsoft JhengHei)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# MTK 風格顏色定義
# =============================================================================

# 背景色
BG_COLOR = RGBColor(255, 249, 230)  # 米色背景
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)

# 強調色（用於區塊標題）
ACCENT_BLUE = RGBColor(70, 130, 180)      # 技術背景、成功判定
ACCENT_ORANGE = RGBColor(230, 126, 34)    # 問題、POC 設計
ACCENT_GREEN = RGBColor(39, 174, 96)      # 效益、解決方案
ACCENT_PURPLE = RGBColor(142, 68, 173)    # 架構、技術細節
ACCENT_RED = RGBColor(192, 0, 0)          # 行動、決策

# 字體
FONT_NAME = "Microsoft JhengHei"


# =============================================================================
# 基本設定
# =============================================================================

def create_presentation():
    """建立 16:9 簡報"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs):
    """加入空白投影片"""
    blank_layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank_layout)


def add_background(slide, prs, color=BG_COLOR):
    """加入背景色"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    return background


# =============================================================================
# MTK 風格：Content Box（圓角矩形 + 標題 + 內容）
# =============================================================================

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    """
    加入 MTK 風格的內容區塊（圓角矩形）

    參數：
    - left, top, width, height: 位置與大小（英吋）
    - title: 區塊標題
    - content_lines: 內容行（list of str）
    - title_color: 標題顏色
    """
    # 圓角矩形背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.03),
        Inches(width - 0.16), Inches(0.3)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT_NAME

    # 內容
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.28),
        Inches(width - 0.16), Inches(height - 0.32)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_NAME
        p.space_after = Pt(1)

    return shape


# =============================================================================
# MTK 風格：標題區
# =============================================================================

def add_main_title(slide, title, subtitle=None):
    """
    加入主標題（含副標題）

    參數：
    - title: 主標題文字
    - subtitle: 副標題文字（可選）
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.25), Inches(0.1), Inches(10), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.5), Inches(12), Inches(0.3)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.font.name = FONT_NAME


# =============================================================================
# MTK 風格：表格
# =============================================================================

def set_cell_text(cell, text, font_size=10, bold=True, color=DARK_GRAY):
    """設定表格儲存格文字"""
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = FONT_NAME


def add_table(slide, left, top, width, height, headers, data, header_color=ACCENT_BLUE):
    """
    加入表格

    參數：
    - headers: 表頭列表 ['欄1', '欄2', '欄3']
    - data: 資料列表 [['A1', 'A2', 'A3'], ['B1', 'B2', 'B3']]
    - header_color: 表頭背景色
    """
    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # 設定表頭
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=10, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = header_color

    # 設定資料
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=9, bold=True)
            # 交替行背景色
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    return table


# =============================================================================
# 基本元件（保留相容性）
# =============================================================================

def add_title(slide, text, left=0.5, top=0.3, width=12.3, height=0.8):
    """加入標題文字框"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_text_block(slide, text, left, top, width, height, font_size=11):
    """加入一般文字區塊"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.LEFT
    return txBox


def add_bullet_list(slide, items, left, top, width, height, title=None, font_size=11):
    """加入項目符號清單"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size + 2)
        p.font.bold = True
        p.font.name = FONT_NAME

    for i, item in enumerate(items):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.name = FONT_NAME
        p.level = 0
    return txBox


def add_image(slide, image_path, left, top, width=None, height=None):
    """加入圖片"""
    kwargs = {'left': Inches(left), 'top': Inches(top)}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(image_path, **kwargs)


def add_section_header(slide, text, left, top, width, bg_color=(0, 112, 192)):
    """加入區塊標題（帶底色的矩形）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    return shape


# =============================================================================
# MTK 風格佈局範例
# =============================================================================

def example_mtk_layout():
    """
    MTK 風格佈局範例：多區塊 + 表格

    ┌──────────────────────────────────────────────────┐
    │  主標題                                           │
    │  副標題                                           │
    ├───────────────┬───────────────┬──────────────────┤
    │  區塊1        │  區塊2        │  區塊3           │
    │  (藍色標題)   │  (橙色標題)   │  (綠色標題)      │
    ├───────────────┴───────────────┼──────────────────┤
    │       表格                     │  區塊4           │
    │                               │  (預期效益)      │
    ├───────────────────────────────┼──────────────────┤
    │  區塊5 (POC 設計)             │  區塊6           │
    │                               │  (成功判定)      │
    └───────────────────────────────┴──────────────────┘
    """
    prs = create_presentation()
    slide = add_blank_slide(prs)

    # 背景
    add_background(slide, prs)

    # 主標題
    add_main_title(
        slide,
        "Cache Aware Scheduling - 解決掉幀問題",
        "靈感來源：Linux sched_cache (Intel/AMD) | 目標：降低跨 cluster 遷移，改善效能穩定度與功耗"
    )

    # 上排三個區塊
    add_content_box(
        slide, 0.25, 0.85, 4.25, 2.0,
        "PC sched_cache 成功要素",
        [
            "在 AMD EPYC Genoa 上驗證有效的前提：",
            "",
            "• 平台具備多個 cache domain (NUMA/CCX)",
            "• Thread 屬於長壽命、持續執行型",
            "• Thread 之間存在高度資料共享",
            "• 跨 cache domain 遷移有實質 cache 破壞成本",
            "",
            "效果：Host time -45%，Throughput +82%",
        ],
        title_color=ACCENT_BLUE
    )

    add_content_box(
        slide, 4.55, 0.85, 4.25, 2.0,
        "Dimensity 9500 Cache 架構",
        [
            "CPU 架構：1×Ultra + 3×Big + 4×Little",
            "",
            "Cache 層級：",
            "• 每個 cluster 內共用一顆 L2 cache",
            "• 不同 cluster 之間不共享 L2",
            "",
            "SoC 層級快取：",
            "• L3 cache：16MB / SLC：10MB（跨 cluster 共用）",
        ],
        title_color=ACCENT_PURPLE
    )

    add_content_box(
        slide, 8.85, 0.85, 4.2, 2.0,
        "為何可能有類似效果",
        [
            "若以下條件在手機遊戲場景成立：",
            "",
            "• GameThread/RenderThread 長壽命、持續執行",
            "• Threads 間有穩定資料共享 (scene/render state)",
            "• 跨 cluster 遷移導致 L2 內容不可沿用",
            "",
            "→ 降低跨 cluster migration 有機會改善穩定度",
        ],
        title_color=ACCENT_GREEN
    )

    # 中間：表格
    add_table(
        slide, 0.25, 3.0, 6.5, 1.0,
        ["平台", "Cache Domain 邊界", "特性"],
        [
            ["AMD Genoa", "CCX/NUMA (L2不共享)", "跨 CCX 遷移成本高"],
            ["Dimensity", "Cluster (L2不共享)", "跨 Cluster 遷移在 L2 層級即有成本"],
        ],
        header_color=ACCENT_BLUE
    )

    # 中間右：預期效益
    add_content_box(
        slide, 6.85, 2.95, 6.2, 1.7,
        "預期效益",
        [
            "若能實作類似 Cache-Aware 的排程機制：",
            "",
            "• 消除無效的跨 cluster 遷移",
            "• 確保 GameThread/RenderThread 留在同 cluster",
            "• 收斂 Frame Time 抖動（1% Low 提升）",
            "• 功耗行為改善（不需急拉頻率救 delay）",
        ],
        title_color=ACCENT_GREEN
    )

    # 下排：POC 設計 + 成功判定
    add_content_box(
        slide, 0.25, 4.7, 6.5, 2.55,
        "POC 設計",
        [
            "實驗條件：",
            "• A：現行方案 (FPSGO/Loom)",
            "• B：現行方案 (Loom) + LAVD (cache aware)",
            "",
            "遊戲場景：",
            "① 遊戲高刷（穩態）",
            "② 遊戲 + 抖音複合場景（干擾/搶佔）",
            "",
            "觀測指標：",
            "• 關鍵 thread 的 cluster 遷移次數",
            "• L2 cache miss rate 變化",
        ],
        title_color=ACCENT_ORANGE
    )

    add_content_box(
        slide, 6.85, 4.7, 6.2, 2.55,
        "成功判定準則",
        [
            "在 Avg FPS 相近（或 target FPS 相同）前提下：",
            "",
            "1. 1% Low 明顯提升",
            "   Frame Time 穩定度改善",
            "",
            "2. 關鍵 Thread 延遲下降",
            "   wakeup latency / runnable delay 顯著下降",
            "",
            "3. 功耗波動不惡化",
            "   理想是下降；至少不因救 delay 而急拉頻率",
        ],
        title_color=ACCENT_BLUE
    )

    return prs


# =============================================================================
# 儲存
# =============================================================================

def save_presentation(prs, output_path):
    """儲存簡報"""
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


# =============================================================================
# 測試
# =============================================================================

if __name__ == '__main__':
    # 測試 MTK 風格佈局
    prs = example_mtk_layout()
    save_presentation(prs, 'test_mtk_style.pptx')
    print("已生成 MTK 風格範例")
