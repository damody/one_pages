"""折線圖（原生圖表）"""
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ._colors import FONT_NAME


def draw_line_chart(slide, left, top, width, height, title, categories, series,
                    show_legend=True, show_data_labels=False, smooth=False):
    """
    繪製折線圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        categories: X 軸類別 ["Week 1", "Week 2", ...]
        series: 資料序列 [{"name": "延遲", "values": [80, 65, 45, 30], "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_data_labels: 是否顯示資料標籤
        smooth: 是否平滑曲線
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series:
        chart_data.add_series(s["name"], s["values"])

    chart_type = XL_CHART_TYPE.LINE_MARKERS_STACKED if not smooth else XL_CHART_TYPE.LINE_MARKERS

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    for i, s in enumerate(series):
        if "color" in s and i < len(chart.series):
            chart.series[i].format.line.color.rgb = s["color"]

    return graphic_frame
