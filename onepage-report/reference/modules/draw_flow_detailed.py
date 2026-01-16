"""詳細版橫向流程圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_RED, COLOR_BLUE, COLOR_TEXT, COLOR_WHITE, FONT_NAME


def draw_flow_detailed(slide, left, top, width, height, nodes, arrow_labels=None, show_highlight=True):
    """
    繪製詳細版橫向流程圖，支援箭頭上的文字標籤和高亮節點

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        nodes: 節點列表，每個元素是 {
            "title": "節點標題",
            "desc": "說明文字（可選）",
            "time": "時間標籤（可選）",
            "color": 顏色（可選，預設 COLOR_BLUE）,
            "highlight": True/False（是否高亮，用紅色虛線框）
        }
        arrow_labels: 箭頭上的文字標籤列表（長度應為 len(nodes)-1）
        show_highlight: 是否顯示高亮標記
    """
    node_count = len(nodes)
    gap = 0.18
    arrow_width = 0.15
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            time_label = node.get("time", "")
            color = node.get("color", COLOR_BLUE)
            highlight = node.get("highlight", False)
        else:
            title = str(node)
            desc = ""
            time_label = ""
            color = COLOR_BLUE
            highlight = False

        # 節點矩形
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # 高亮標記（紅色虛線框）
        if highlight and show_highlight:
            highlight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x - 0.03), Inches(top - 0.03),
                Inches(node_width + 0.06), Inches(height + 0.06)
            )
            highlight_box.fill.background()
            highlight_box.line.color.rgb = COLOR_RED
            highlight_box.line.width = Pt(2)
            highlight_box.line.dash_style = 2

        # 節點文字
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(7)
            p2.font.color.rgb = COLOR_WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if time_label:
            p3 = tf.add_paragraph()
            p3.text = time_label
            p3.font.size = Pt(7)
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(255, 255, 200)
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER

        # 箭頭
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.02), Inches(top + height/2 - 0.08),
                Inches(arrow_width), Inches(0.16)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.fill.background()

            if arrow_labels and i < len(arrow_labels) and arrow_labels[i]:
                label_box = slide.shapes.add_textbox(
                    Inches(x + node_width + 0.02), Inches(top - 0.18),
                    Inches(gap - 0.04), Inches(0.18)
                )
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = arrow_labels[i]
                p.font.size = Pt(6)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER
