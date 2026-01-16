"""
輔助函數模組
包含：區塊標題、項目列表、內容方塊
"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_BLUE, COLOR_GRAY_BG, COLOR_TEXT, FONT_NAME


def add_section_title(slide, left, top, width, text, color=COLOR_BLUE):
    """
    加入區塊標題

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width: 寬度（吋）
        text: 標題文字
        color: 標題顏色
    """
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    return box


def add_bullet_list(slide, left, top, width, height, items, font_size=10):
    """
    加入項目列表

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        items: 項目列表
        font_size: 字體大小
    """
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.space_after = Pt(4)

    return box


def add_content_box(slide, left, top, width, height, title, content, title_color=COLOR_BLUE, bg_color=None):
    """
    加入帶標題的內容方塊

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 標題文字
        content: 內容（字串或列表）
        title_color: 標題顏色
        bg_color: 背景顏色（可選）
    """
    if bg_color:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(width - 0.2), Inches(0.3)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT_NAME

    # 內容
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(width - 0.2), Inches(height - 0.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    if isinstance(content, list):
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_NAME
    else:
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
