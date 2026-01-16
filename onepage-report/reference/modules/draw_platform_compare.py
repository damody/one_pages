"""平台對比圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_GREEN, COLOR_BLUE, COLOR_ORANGE, COLOR_GRAY_BG, COLOR_TEXT, COLOR_WHITE, FONT_NAME
from .draw_flow_detailed import draw_flow_detailed


def draw_platform_compare(slide, left, top, width, height,
                          platform_a, platform_b, differences=None):
    """
    繪製上下平台對比圖，每個平台內部有完整流程

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        platform_a: 上方平台 {
            "name": "PC 平台",
            "title": "AMD Anti-Lag 2 流程（已驗證有效）",
            "color": COLOR_BLUE,
            "flow_nodes": [...],
            "arrow_labels": [...],
            "summary": "效果：CS2 延遲降低 37%"
        }
        platform_b: 下方平台（同上格式）
        differences: 差異標註列表 [{"item": "輸入方式", "a": "滑鼠", "b": "觸控"}, ...]
    """
    diff_width = 0 if not differences else 2.5
    platform_width = width - diff_width - 0.1
    platform_height = (height - 0.2) / 2

    platforms = [platform_a, platform_b]
    colors = [COLOR_BLUE, COLOR_GREEN]

    for i, platform in enumerate(platforms):
        y = top + i * (platform_height + 0.2)
        color = platform.get("color", colors[i])
        name = platform.get("name", f"平台 {i+1}")
        title = platform.get("title", "")
        flow_nodes = platform.get("flow_nodes", [])
        arrow_labels = platform.get("arrow_labels", [])
        summary = platform.get("summary", "")

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(platform_width), Inches(platform_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_GRAY_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)

        name_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.1), Inches(y + 0.08),
            Inches(1.0), Inches(0.28)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = color
        name_box.line.fill.background()

        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(
            Inches(left + 1.2), Inches(y + 0.08),
            Inches(platform_width - 1.4), Inches(0.28)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = FONT_NAME

        if flow_nodes:
            flow_height = platform_height - 0.7 if summary else platform_height - 0.45
            draw_flow_detailed(
                slide, left=left + 0.1, top=y + 0.4,
                width=platform_width - 0.2, height=flow_height,
                nodes=flow_nodes, arrow_labels=arrow_labels, show_highlight=True
            )

        if summary:
            summary_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(y + platform_height - 0.25),
                Inches(platform_width - 0.2), Inches(0.2)
            )
            tf = summary_box.text_frame
            p = tf.paragraphs[0]
            p.text = summary
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = color
            p.font.name = FONT_NAME

    if differences:
        diff_left = left + platform_width + 0.1
        diff_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(diff_left), Inches(top),
            Inches(diff_width), Inches(height)
        )
        diff_box.fill.solid()
        diff_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
        diff_box.line.color.rgb = COLOR_ORANGE
        diff_box.line.width = Pt(1)

        diff_title = slide.shapes.add_textbox(
            Inches(diff_left + 0.1), Inches(top + 0.08),
            Inches(diff_width - 0.2), Inches(0.25)
        )
        tf = diff_title.text_frame
        p = tf.paragraphs[0]
        p.text = "主要差異"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_ORANGE
        p.font.name = FONT_NAME

        item_height = (height - 0.4) / len(differences)
        for j, diff in enumerate(differences):
            item_y = top + 0.35 + j * item_height
            item_box = slide.shapes.add_textbox(
                Inches(diff_left + 0.1), Inches(item_y),
                Inches(diff_width - 0.2), Inches(item_height)
            )
            tf = item_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = diff.get("item", "")
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_NAME

            p2 = tf.add_paragraph()
            p2.text = f"PC: {diff.get('a', '')}"
            p2.font.size = Pt(7)
            p2.font.color.rgb = COLOR_BLUE
            p2.font.name = FONT_NAME

            p3 = tf.add_paragraph()
            p3.text = f"手機: {diff.get('b', '')}"
            p3.font.size = Pt(7)
            p3.font.color.rgb = COLOR_GREEN
            p3.font.name = FONT_NAME
