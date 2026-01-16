"""圓餅圖（原生圖表）"""
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ._colors import FONT_NAME


def draw_pie_chart(slide, left, top, width, height, title, data,
                   show_legend=True, show_percentage=True):
    """
    繪製圓餅圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        data: 資料 [{"name": "項目", "value": 45, "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_percentage: 是否顯示百分比
    """
    chart_data = CategoryChartData()
    chart_data.categories = [d["name"] for d in data]
    chart_data.add_series("Values", [d["value"] for d in data])

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    if len(chart.series) > 0:
        for i, d in enumerate(data):
            if "color" in d and i < len(chart.series[0].points):
                chart.series[0].points[i].format.fill.solid()
                chart.series[0].points[i].format.fill.fore_color.rgb = d["color"]

    return graphic_frame
