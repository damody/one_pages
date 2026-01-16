"""帶內部流程圖的前後對比"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_RED, COLOR_GREEN, COLOR_BLUE, COLOR_GRAY_BG, FONT_NAME
from .draw_flow_detailed import draw_flow_detailed
from .draw_comparison_table import draw_comparison_table


def draw_before_after_with_flow(slide, left, top, width, height,
                                 before_title, before_flow_nodes, before_arrow_labels,
                                 after_title, after_flow_nodes, after_arrow_labels,
                                 center_arrow_label="導入方案", bottom_table=None):
    """
    繪製帶內部流程圖的前後對比

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        before_title: 左側標題（如「改善前：畫面堆積導致延遲」）
        before_flow_nodes: 左側內部流程節點列表
        before_arrow_labels: 左側箭頭標籤列表
        after_title: 右側標題
        after_flow_nodes: 右側內部流程節點列表
        after_arrow_labels: 右側箭頭標籤列表
        center_arrow_label: 中間箭頭上的文字（如「導入 SDK」）
        bottom_table: 底部對比表格 {"headers": [...], "rows": [[...], ...]}
    """
    from pptx.enum.text import PP_ALIGN

    table_height = 0.7 if bottom_table else 0
    main_height = height - table_height
    box_width = (width - 0.5) / 2
    flow_area_height = main_height - 0.5

    # 左側區塊（改善前）
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(main_height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.35)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    if before_flow_nodes:
        draw_flow_detailed(
            slide, left=left + 0.1, top=top + 0.45,
            width=box_width - 0.2, height=flow_area_height - 0.1,
            nodes=before_flow_nodes, arrow_labels=before_arrow_labels,
            show_highlight=True
        )

    # 右側區塊（改善後）
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.5), Inches(top),
        Inches(box_width), Inches(main_height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    after_title_box = slide.shapes.add_textbox(
        Inches(left + box_width + 0.6), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.35)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    if after_flow_nodes:
        draw_flow_detailed(
            slide, left=left + box_width + 0.6, top=top + 0.45,
            width=box_width - 0.2, height=flow_area_height - 0.1,
            nodes=after_flow_nodes, arrow_labels=after_arrow_labels,
            show_highlight=False
        )

    # 中間箭頭
    arrow_y = top + main_height / 2 - 0.15
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.1), Inches(arrow_y),
        Inches(0.3), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()

    if center_arrow_label:
        center_label = slide.shapes.add_textbox(
            Inches(left + box_width + 0.05), Inches(arrow_y - 0.22),
            Inches(0.4), Inches(0.2)
        )
        tf = center_label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = center_arrow_label
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLUE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    if bottom_table:
        headers = bottom_table.get("headers", [])
        rows = bottom_table.get("rows", [])
        if headers and rows:
            draw_comparison_table(
                slide, left=left, top=top + main_height + 0.05,
                width=width, height=table_height - 0.05,
                headers=headers, rows=rows
            )
