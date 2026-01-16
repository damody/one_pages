"""對比表格"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_BLUE, COLOR_GRAY_BG, COLOR_TEXT, COLOR_WHITE, FONT_NAME


def draw_comparison_table(slide, left, top, width, height, headers, rows):
    """
    繪製對比表格

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        headers: ["項目", "PC", "手機"]
        rows: [["延遲", "20ms", "50ms"], ...]
    """
    col_count = len(headers)
    row_count = len(rows) + 1  # +1 for header
    col_width = width / col_count
    row_height = height / row_count

    for r in range(row_count):
        for c in range(col_count):
            x = left + c * col_width
            y = top + r * row_height

            # 儲存格背景
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(col_width), Inches(row_height)
            )

            if r == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BLUE
                text_color = COLOR_WHITE
                is_bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_GRAY_BG if r % 2 == 1 else COLOR_WHITE
                text_color = COLOR_TEXT
                is_bold = False

            cell.line.color.rgb = RGBColor(200, 200, 200)
            cell.line.width = Pt(0.5)

            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(y + 0.05),
                Inches(col_width - 0.1), Inches(row_height - 0.1)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            if r == 0:
                p.text = headers[c]
            else:
                p.text = rows[r-1][c]

            p.font.size = Pt(9)
            p.font.bold = is_bold
            p.font.color.rgb = text_color
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
