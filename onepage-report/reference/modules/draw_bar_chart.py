"""長條圖（原生圖表）"""
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ._colors import FONT_NAME


def draw_bar_chart(slide, left, top, width, height, title, categories, series,
                   show_legend=True, show_data_labels=False, horizontal=False):
    """
    繪製長條圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        categories: X 軸類別 ["項目A", "項目B", ...]
        series: 資料序列 [{"name": "數據", "values": [10, 20, 30], "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_data_labels: 是否顯示資料標籤
        horizontal: 是否水平顯示
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series:
        chart_data.add_series(s["name"], s["values"])

    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    for i, s in enumerate(series):
        if "color" in s and i < len(chart.series):
            for point in chart.series[i].points:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = s["color"]

    return graphic_frame
