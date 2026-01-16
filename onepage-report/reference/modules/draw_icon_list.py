"""帶圖標的列表"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_RED, COLOR_GREEN, COLOR_BLUE, COLOR_ORANGE, COLOR_TEXT, FONT_NAME


def draw_icon_list(slide, left, top, width, item_height, items):
    """
    繪製帶圖標的列表

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width: 寬度（吋）
        item_height: 每個項目的高度（吋）
        items: [{"icon": "check/cross/warn", "text": "...", "color": ...}, ...]
    """
    for i, item in enumerate(items):
        y = top + i * item_height
        icon_type = item.get("icon", "check")
        color = item.get("color", COLOR_GREEN if icon_type == "check" else COLOR_RED)

        # 圖標
        icon_size = 0.2
        if icon_type == "check":
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_GREEN
            circle.line.fill.background()
        elif icon_type == "cross":
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_RED
            circle.line.fill.background()
        elif icon_type == "warn":
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            triangle.fill.solid()
            triangle.fill.fore_color.rgb = COLOR_ORANGE
            triangle.line.fill.background()
        else:
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_BLUE
            circle.line.fill.background()

        # 文字
        text_box = slide.shapes.add_textbox(
            Inches(left + icon_size + 0.1), Inches(y),
            Inches(width - icon_size - 0.1), Inches(item_height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("text", "")
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
