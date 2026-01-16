"""迷你流程圖（用於術語卡片）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_BLUE, COLOR_WHITE, FONT_NAME


def draw_mini_flow(slide, left, top, width, height, nodes):
    """
    繪製迷你流程圖（用於術語卡片內的示意圖）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        nodes: 節點列表，可以是字串或 dict {"text": "...", "color": ...}
    """
    node_count = len(nodes)
    gap = 0.08
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)
        color = node.get("color", COLOR_BLUE) if isinstance(node, dict) else COLOR_BLUE
        text = node.get("text", node) if isinstance(node, dict) else str(node)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.01), Inches(top + height/2 - 0.04),
                Inches(gap - 0.02), Inches(0.08)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()
