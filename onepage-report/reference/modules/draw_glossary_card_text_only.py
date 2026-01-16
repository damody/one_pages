"""純文字術語卡片"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_BLUE, COLOR_GRAY_BG, COLOR_TEXT, FONT_NAME


def draw_glossary_card_text_only(slide, left, top, width, height, term, desc):
    """
    繪製純文字術語卡片

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        term: 術語名稱
        desc: 術語解釋（簡短版）
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_BG
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(0.5)

    term_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.05),
        Inches(width - 0.16), Inches(0.28)
    )
    tf = term_box.text_frame
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.33),
        Inches(width - 0.16), Inches(height - 0.4)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
