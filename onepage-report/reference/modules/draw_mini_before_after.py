"""迷你前後對比圖（用於術語卡片）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_RED, COLOR_GREEN, COLOR_BLUE, FONT_NAME


def draw_mini_before_after(slide, left, top, width, height, before_text, after_text):
    """
    繪製迷你前後對比圖（用於術語卡片內的示意圖）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        before_text: 左側（改善前）文字
        after_text: 右側（改善後）文字
    """
    box_width = (width - 0.25) / 2

    # 左側（Before）
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(1)

    tf = before_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = before_text
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 右側（After）
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.25), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(1)

    tf = after_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = after_text
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 中間箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.05), Inches(top + height/2 - 0.06),
        Inches(0.15), Inches(0.12)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()
