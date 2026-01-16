"""前後對比圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_RED, COLOR_GREEN, COLOR_BLUE, COLOR_GRAY_BG, COLOR_TEXT, FONT_NAME


def draw_before_after(slide, left, top, width, height, before_title, before_items, after_title, after_items):
    """
    繪製前後對比圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        before_title: 左側標題（如「現況問題」）
        before_items: 左側項目列表
        after_title: 右側標題（如「導入 SDK 後」）
        after_items: 右側項目列表
    """
    box_width = (width - 0.4) / 2

    # 左側區塊
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    before_content = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(box_width - 0.2), Inches(height - 0.5)
    )
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(before_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

    # 右側區塊
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.4), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    after_title_box = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    after_content = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.4),
        Inches(box_width - 0.2), Inches(height - 0.5)
    )
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(after_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

    # 中間箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.1), Inches(top + height/2 - 0.12),
        Inches(0.2), Inches(0.24)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()
