"""分層架構圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ._colors import COLOR_BLUE, COLOR_TEXT, COLOR_WHITE, FONT_NAME


def draw_architecture(slide, left, top, width, height, layers):
    """
    繪製分層架構圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        layers: [{"name": "...", "color": ..., "components": ["...", ...]}, ...]
    """
    layer_count = len(layers)
    layer_height = (height - 0.1 * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + 0.1)
        color = layer.get("color", COLOR_BLUE)

        # 層背景
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # 層名稱（左側）
        name_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(y + layer_height/2 - 0.15),
            Inches(1.2), Inches(0.3)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer.get("name", "")
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        # 組件（右側橫排）
        components = layer.get("components", [])
        if components:
            comp_width = (width - 1.5) / len(components)
            for j, comp in enumerate(components):
                cx = left + 1.4 + j * comp_width
                comp_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(y + 0.1),
                    Inches(comp_width - 0.1), Inches(layer_height - 0.2)
                )
                comp_box.fill.solid()
                comp_box.fill.fore_color.rgb = COLOR_WHITE
                comp_box.line.fill.background()

                tf = comp_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = comp
                p.font.size = Pt(8)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER
