"""迷你時間軸（用於術語卡片）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_BLUE, COLOR_WHITE, FONT_NAME


def draw_mini_timeline(slide, left, top, width, height, stages):
    """
    繪製迷你時間軸（用於術語卡片內的示意圖）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        stages: 階段列表，可以是字串或 dict {"text": "...", "color": ...}
    """
    line_y = top + height * 0.6
    stage_count = len(stages)
    stage_width = width / stage_count

    # 水平軸線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(line_y),
        Inches(width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(100, 100, 100)
    line.line.fill.background()

    for i, stage in enumerate(stages):
        x = left + i * stage_width
        color = stage.get("color", COLOR_BLUE) if isinstance(stage, dict) else COLOR_BLUE
        text = stage.get("text", stage) if isinstance(stage, dict) else str(stage)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.02), Inches(top),
            Inches(stage_width - 0.04), Inches(height * 0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(6)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
