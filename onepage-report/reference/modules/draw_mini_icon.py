"""迷你圖標（用於術語卡片）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ._colors import COLOR_BLUE, COLOR_GREEN, COLOR_ORANGE, COLOR_TEXT, FONT_NAME


def draw_mini_icon(slide, left, top, width, height, icon_type, label):
    """
    繪製迷你圖標（用於術語卡片內的示意圖）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        icon_type: 圖標類型 ("chart", "gauge", "warning", 或其他預設圓形)
        label: 標籤文字
    """
    icon_size = min(width, height * 0.6)
    icon_x = left + (width - icon_size) / 2
    icon_y = top

    if icon_type == "chart":
        bar_width = icon_size / 4
        bars = [0.4, 0.7, 0.5, 0.9]
        for i, h in enumerate(bars):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(icon_x + i * bar_width), Inches(icon_y + icon_size * (1 - h)),
                Inches(bar_width - 0.02), Inches(icon_size * h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_BLUE
            bar.line.fill.background()
    elif icon_type == "gauge":
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_GREEN
        circle.line.fill.background()
    elif icon_type == "warning":
        triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = COLOR_ORANGE
        triangle.line.fill.background()
    else:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_BLUE
        circle.line.fill.background()

    label_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + icon_size + 0.05),
        Inches(width), Inches(height - icon_size - 0.05)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
