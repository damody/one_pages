"""指標卡片"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ._colors import COLOR_BLUE, COLOR_WHITE, FONT_NAME


def draw_metric_cards(slide, left, top, width, height, metrics):
    """
    繪製指標卡片（數字 + 說明）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        metrics: [{"value": "81%", "label": "延遲降低", "color": ...}, ...]
    """
    card_count = len(metrics)
    card_width = (width - 0.15 * (card_count - 1)) / card_count

    for i, metric in enumerate(metrics):
        x = left + i * (card_width + 0.15)
        color = metric.get("color", COLOR_BLUE)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(card_width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        # 數值（大字）
        value_box = slide.shapes.add_textbox(
            Inches(x), Inches(top + 0.1),
            Inches(card_width), Inches(height * 0.5)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get("value", "")
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 標籤（小字）
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(top + height * 0.55),
            Inches(card_width), Inches(height * 0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get("label", "")
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
