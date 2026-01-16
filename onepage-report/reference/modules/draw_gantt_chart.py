"""甘特圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta

from ._colors import COLOR_RED, COLOR_BLUE, COLOR_GRAY_DARK, COLOR_TEXT, FONT_NAME


def draw_gantt_chart(slide, left, top, width, height, title, tasks, milestones=None,
                     time_unit="week", show_today_line=True, show_progress=True):
    """
    繪製甘特圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        tasks: [{"name": "任務", "start": "2025-01-06", "end": "2025-01-19",
                 "progress": 100, "color": COLOR_BLUE}, ...]
        milestones: [{"name": "里程碑", "date": "2025-02-28", "color": COLOR_RED}, ...]
        time_unit: "day" | "week" | "month"
        show_today_line: 是否顯示今日線
        show_progress: 是否顯示進度條
    """
    def parse_date(date_str):
        return datetime.strptime(date_str, "%Y-%m-%d")

    all_dates = []
    for t in tasks:
        all_dates.append(parse_date(t["start"]))
        all_dates.append(parse_date(t["end"]))
    if milestones:
        for m in milestones:
            all_dates.append(parse_date(m["date"]))

    min_date = min(all_dates)
    max_date = max(all_dates)
    total_days = (max_date - min_date).days + 1

    title_height = 0.35
    header_height = 0.3
    name_col_width = 1.8
    chart_left = left + name_col_width
    chart_width = width - name_col_width
    task_count = len(tasks)
    task_height = (height - title_height - header_height) / task_count if task_count > 0 else 0.5

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(title_height)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    header_top = top + title_height

    # 計算時間刻度
    if time_unit == "week":
        current = min_date
        ticks = []
        while current <= max_date:
            ticks.append(current)
            current += timedelta(days=7)
        if ticks[-1] < max_date:
            ticks.append(max_date)
    else:
        tick_count = min(8, total_days)
        ticks = [min_date + timedelta(days=i * total_days // tick_count) for i in range(tick_count + 1)]

    # 繪製時間刻度
    for i, tick_date in enumerate(ticks[:-1]):
        tick_x = chart_left + (tick_date - min_date).days / total_days * chart_width
        next_tick_x = chart_left + (ticks[i + 1] - min_date).days / total_days * chart_width if i + 1 < len(ticks) else chart_left + chart_width

        label_box = slide.shapes.add_textbox(
            Inches(tick_x), Inches(header_top),
            Inches(next_tick_x - tick_x), Inches(header_height)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = tick_date.strftime("%m/%d")
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(tick_x), Inches(header_top + header_height),
            Inches(0.01), Inches(height - title_height - header_height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(220, 220, 220)
        line.line.fill.background()

    # 繪製任務
    task_area_top = header_top + header_height
    for i, task in enumerate(tasks):
        task_top = task_area_top + i * task_height
        task_start = parse_date(task["start"])
        task_end = parse_date(task["end"])
        color = task.get("color", COLOR_BLUE)
        progress = task.get("progress", 0)

        name_box = slide.shapes.add_textbox(
            Inches(left), Inches(task_top),
            Inches(name_col_width - 0.1), Inches(task_height)
        )
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = task["name"]
        p.font.size = Pt(8)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

        start_offset = (task_start - min_date).days / total_days
        end_offset = (task_end - min_date).days / total_days
        bar_left = chart_left + start_offset * chart_width
        bar_width = (end_offset - start_offset) * chart_width
        bar_top = task_top + task_height * 0.2
        bar_height = task_height * 0.6

        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_left), Inches(bar_top),
            Inches(bar_width), Inches(bar_height)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(230, 230, 230)
        bg_bar.line.fill.background()

        if show_progress and progress > 0:
            progress_width = bar_width * progress / 100
            progress_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(bar_top),
                Inches(progress_width), Inches(bar_height)
            )
            progress_bar.fill.solid()
            progress_bar.fill.fore_color.rgb = color
            progress_bar.line.fill.background()

    # 繪製里程碑
    if milestones:
        for m in milestones:
            m_date = parse_date(m["date"])
            m_offset = (m_date - min_date).days / total_days
            m_x = chart_left + m_offset * chart_width
            m_color = m.get("color", COLOR_RED)

            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(m_x - 0.1), Inches(task_area_top - 0.05),
                Inches(0.2), Inches(0.2)
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = m_color
            diamond.line.fill.background()

    # 今日線
    if show_today_line:
        today = datetime.now()
        if min_date <= today <= max_date:
            today_offset = (today - min_date).days / total_days
            today_x = chart_left + today_offset * chart_width
            today_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(today_x), Inches(header_top),
                Inches(0.02), Inches(height - title_height)
            )
            today_line.fill.solid()
            today_line.fill.fore_color.rgb = COLOR_RED
            today_line.line.fill.background()
