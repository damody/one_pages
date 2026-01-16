"""橫向流程圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_BLUE, COLOR_WHITE, FONT_NAME


def draw_flow(slide, left, top, width, height, nodes):
    """
    繪製橫向流程圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        nodes: 節點列表，可以是字串或 dict
            字串: 直接作為標題
            dict: {"title": "...", "desc": "...", "color": ...}
    """
    node_count = len(nodes)
    gap = 0.12
    arrow_width = 0.1
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            color = node.get("color", COLOR_BLUE)
        else:
            title = str(node)
            desc = ""
            color = COLOR_BLUE

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(8)
            p2.font.color.rgb = COLOR_WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.01), Inches(top + height/2 - 0.06),
                Inches(arrow_width), Inches(0.12)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()
