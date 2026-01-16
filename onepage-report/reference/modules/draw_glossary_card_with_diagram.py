"""帶示意圖的術語卡片"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ._colors import COLOR_BLUE, COLOR_GRAY_BG, COLOR_TEXT, FONT_NAME
from .draw_mini_flow import draw_mini_flow
from .draw_mini_before_after import draw_mini_before_after
from .draw_mini_layers import draw_mini_layers
from .draw_mini_timeline import draw_mini_timeline
from .draw_mini_icon import draw_mini_icon


def draw_glossary_card_with_diagram(slide, left, top, width, height, term, desc, diagram_type, diagram_params):
    """
    繪製帶示意圖的術語卡片

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        term: 術語名稱
        desc: 術語解釋
        diagram_type: "flow" | "before_after" | "layers" | "timeline" | "icon"
        diagram_params: 示意圖參數
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_BG
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(1)

    diagram_height = height * 0.45
    diagram_top = top + 0.1
    diagram_left = left + 0.15
    diagram_width = width - 0.3

    if diagram_type == "flow":
        draw_mini_flow(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("nodes", []))
    elif diagram_type == "before_after":
        draw_mini_before_after(slide, diagram_left, diagram_top, diagram_width, diagram_height,
                               diagram_params.get("before", ""), diagram_params.get("after", ""))
    elif diagram_type == "layers":
        draw_mini_layers(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("layers", []))
    elif diagram_type == "timeline":
        draw_mini_timeline(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("stages", []))
    elif diagram_type == "icon":
        draw_mini_icon(slide, diagram_left, diagram_top, diagram_width, diagram_height,
                       diagram_params.get("icon_type", ""), diagram_params.get("label", ""))

    term_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + diagram_height + 0.15),
        Inches(width - 0.2), Inches(0.35)
    )
    tf = term_box.text_frame
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + diagram_height + 0.5),
        Inches(width - 0.2), Inches(height - diagram_height - 0.6)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
