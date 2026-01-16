"""迷你分層圖（用於術語卡片）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ._colors import COLOR_BLUE, COLOR_WHITE, FONT_NAME


def draw_mini_layers(slide, left, top, width, height, layers):
    """
    繪製迷你分層圖（用於術語卡片內的示意圖）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        layers: 層列表，可以是字串或 dict {"text": "...", "color": ...}
    """
    layer_count = len(layers)
    layer_height = (height - 0.05 * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + 0.05)
        color = layer.get("color", COLOR_BLUE) if isinstance(layer, dict) else COLOR_BLUE
        text = layer.get("text", layer) if isinstance(layer, dict) else str(layer)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
