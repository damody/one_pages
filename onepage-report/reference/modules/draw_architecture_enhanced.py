"""增強版系統架構圖"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ._colors import COLOR_BLUE, COLOR_ORANGE, COLOR_GRAY_BG, COLOR_GRAY_DARK, COLOR_TEXT, COLOR_WHITE, FONT_NAME


def draw_architecture_enhanced(slide, left, top, width, height, layers,
                                connections=None, annotations=None):
    """
    繪製增強版系統架構圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        layers: [{"name": "層名", "color": ..., "components": [
            {"name": "組件", "highlight": False, "highlight_label": ""}, ...]}, ...]
        connections: [{"from": "組件A", "to": "組件B", "label": "...", "style": "arrow"}, ...]
        annotations: [{"target": "組件", "text": "說明", "position": "right"}, ...]
    """
    component_positions = {}

    layer_count = len(layers)
    layer_gap = 0.15
    layer_height = (height - layer_gap * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + layer_gap)
        color = layer.get("color", COLOR_BLUE)
        name = layer.get("name", "")

        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = COLOR_GRAY_BG
        layer_box.line.color.rgb = color
        layer_box.line.width = Pt(2)

        name_width = 1.2
        name_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.05), Inches(y + 0.05),
            Inches(name_width - 0.1), Inches(layer_height - 0.1)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = color
        name_box.line.fill.background()

        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        components = layer.get("components", [])
        if components:
            comp_area_width = width - name_width - 0.2
            comp_width = comp_area_width / len(components) if len(components) > 0 else 1
            comp_height = layer_height - 0.2

            for j, comp in enumerate(components):
                comp_name = comp["name"] if isinstance(comp, dict) else comp
                highlight = comp.get("highlight", False) if isinstance(comp, dict) else False
                highlight_label = comp.get("highlight_label", "") if isinstance(comp, dict) else ""

                cx = left + name_width + 0.1 + j * comp_width
                cy = y + 0.1

                comp_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(cy),
                    Inches(comp_width - 0.1), Inches(comp_height)
                )
                comp_box.fill.solid()
                comp_box.fill.fore_color.rgb = COLOR_WHITE
                comp_box.line.color.rgb = color
                comp_box.line.width = Pt(1)

                if highlight:
                    comp_box.line.color.rgb = COLOR_ORANGE
                    comp_box.line.width = Pt(3)

                    if highlight_label:
                        label_box = slide.shapes.add_textbox(
                            Inches(cx), Inches(cy - 0.18),
                            Inches(comp_width - 0.1), Inches(0.18)
                        )
                        tf = label_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = highlight_label
                        p.font.size = Pt(6)
                        p.font.bold = True
                        p.font.color.rgb = COLOR_ORANGE
                        p.font.name = FONT_NAME
                        p.alignment = PP_ALIGN.CENTER

                tf = comp_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = comp_name
                p.font.size = Pt(8)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER

                component_positions[comp_name] = {
                    "left": cx, "top": cy,
                    "width": comp_width - 0.1, "height": comp_height,
                    "center_x": cx + (comp_width - 0.1) / 2,
                    "center_y": cy + comp_height / 2,
                    "bottom": cy + comp_height
                }

    if connections:
        for conn in connections:
            from_pos = component_positions.get(conn["from"])
            to_pos = component_positions.get(conn["to"])

            if from_pos and to_pos:
                start_x = from_pos["center_x"]
                start_y = from_pos["bottom"]
                end_x = to_pos["center_x"]
                end_y = to_pos["top"]
                label = conn.get("label", "")

                if abs(start_x - end_x) < 0.1:
                    line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(start_x - 0.01), Inches(start_y),
                        Inches(0.02), Inches(end_y - start_y)
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = COLOR_GRAY_DARK
                    line.line.fill.background()
                else:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(min(start_x, end_x)), Inches(start_y),
                        Inches(abs(end_x - start_x) + 0.1), Inches(end_y - start_y)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = COLOR_GRAY_DARK
                    arrow.line.fill.background()

                if label:
                    label_x = (start_x + end_x) / 2 - 0.3
                    label_y = (start_y + end_y) / 2 - 0.1
                    label_box = slide.shapes.add_textbox(
                        Inches(label_x), Inches(label_y),
                        Inches(0.6), Inches(0.2)
                    )
                    tf = label_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(6)
                    p.font.color.rgb = COLOR_GRAY_DARK
                    p.font.name = FONT_NAME
                    p.alignment = PP_ALIGN.CENTER

    return component_positions
