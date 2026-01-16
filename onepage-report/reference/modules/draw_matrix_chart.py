"""矩陣圖（風險評估/優先級矩陣）"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ._colors import COLOR_ORANGE, COLOR_TEXT, COLOR_WHITE, FONT_NAME


def draw_matrix_chart(slide, left, top, width, height, title, x_axis, y_axis, items, quadrant_labels=None):
    """
    繪製矩陣圖（風險評估/優先級矩陣）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        x_axis: {"label": "影響程度", "values": ["低", "中", "高"]}
        y_axis: {"label": "發生機率", "values": ["低", "中", "高"]}
        items: [{"name": "風險項", "x": 2, "y": 1, "color": COLOR_ORANGE}, ...]
        quadrant_labels: {"high-high": "高優先", "low-low": "低優先", ...}
    """
    title_height = 0.35
    axis_label_width = 0.4
    axis_value_height = 0.25

    chart_left = left + axis_label_width + 0.1
    chart_top = top + title_height
    chart_width = width - axis_label_width - 0.2
    chart_height = height - title_height - axis_value_height - 0.1

    x_count = len(x_axis["values"])
    y_count = len(y_axis["values"])
    cell_width = chart_width / x_count
    cell_height = chart_height / y_count

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(title_height)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # Y 軸標籤
    y_label_box = slide.shapes.add_textbox(
        Inches(left), Inches(chart_top + chart_height / 2 - 0.2),
        Inches(axis_label_width), Inches(0.4)
    )
    tf = y_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = y_axis["label"]
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # X 軸標籤
    x_label_box = slide.shapes.add_textbox(
        Inches(chart_left + chart_width / 2 - 0.4), Inches(top + height - 0.25),
        Inches(0.8), Inches(0.25)
    )
    tf = x_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = x_axis["label"]
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 繪製格子
    for yi, y_val in enumerate(reversed(y_axis["values"])):
        for xi, x_val in enumerate(x_axis["values"]):
            cell_x = chart_left + xi * cell_width
            cell_y = chart_top + yi * cell_height

            risk_level = (xi + (y_count - 1 - yi)) / (x_count + y_count - 2)
            bg_r = int(232 + (255 - 232) * risk_level)
            bg_g = int(245 - (245 - 200) * risk_level)
            bg_b = int(233 - (233 - 200) * risk_level)

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cell_x), Inches(cell_y),
                Inches(cell_width), Inches(cell_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
            cell.line.color.rgb = RGBColor(200, 200, 200)
            cell.line.width = Pt(0.5)

    # Y 軸刻度值
    for yi, y_val in enumerate(reversed(y_axis["values"])):
        val_box = slide.shapes.add_textbox(
            Inches(left + axis_label_width - 0.3), Inches(chart_top + yi * cell_height + cell_height / 2 - 0.1),
            Inches(0.3), Inches(0.2)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = y_val
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.RIGHT

    # X 軸刻度值
    for xi, x_val in enumerate(x_axis["values"]):
        val_box = slide.shapes.add_textbox(
            Inches(chart_left + xi * cell_width), Inches(chart_top + chart_height),
            Inches(cell_width), Inches(axis_value_height)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = x_val
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # 繪製項目點
    for item in items:
        ix = item["x"]
        iy = item["y"]
        color = item.get("color", COLOR_ORANGE)

        point_x = chart_left + (ix + 0.5) * cell_width - 0.15
        point_y = chart_top + (y_count - 1 - iy + 0.5) * cell_height - 0.15

        point = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(point_x), Inches(point_y),
            Inches(0.3), Inches(0.3)
        )
        point.fill.solid()
        point.fill.fore_color.rgb = color
        point.line.color.rgb = COLOR_WHITE
        point.line.width = Pt(2)

        name_box = slide.shapes.add_textbox(
            Inches(point_x - 0.2), Inches(point_y + 0.32),
            Inches(0.7), Inches(0.2)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = item["name"]
        p.font.size = Pt(6)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
