# -*- coding: utf-8 -*-
"""
生成 MTK Latency Meter 一頁式投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 顏色定義
BG_COLOR = RGBColor(255, 249, 230)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_ORANGE = RGBColor(230, 126, 34)

def set_cell_text(cell, text, font_size=10, bold=True, color=DARK_GRAY):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = "Microsoft JhengHei"

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft JhengHei"

    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.16), height - Inches(0.32))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(1)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()

    # ===== 標題區 =====
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MTK Latency Meter - 端到端延遲量測工具"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft JhengHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "靈感來源：AMD Frame Latency Meter / Anti-Lag 2 | 目標：補齊「玩家體感」的可回歸 KPI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 左上區塊：現況痛點 =====
    add_content_box(
        slide, Inches(0.25), Inches(0.85), Inches(4.25), Inches(2.1),
        "現況痛點",
        [
            "現行 KPI：FPS / 1% low / frame time / 功耗 / 溫控",
            "缺少：Input → Display 端到端延遲量測",
            "",
            "現行量測方式問題：",
            "• iPhone 240FPS 高速攝影 + 人工逐幀數幀",
            "• 量測成本高（設備、架設、人工回放）",
            "• 無法大量取樣，難做每日/每週自動回歸",
            "• A/B 迭代慢，跨人/跨場景一致性不足",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 中上區塊：核心價值 =====
    add_content_box(
        slide, Inches(4.55), Inches(0.85), Inches(4.25), Inches(2.1),
        "核心價值",
        [
            "把「高成本人工量測」轉成：",
            "• 可自動化",
            "• 可大量取樣",
            "• 可回歸",
            "",
            "直接回答客戶關心的問題：",
            "「1% low 差不多，為什麼體感不一樣？」",
            "「同功耗下，哪方案 input lag 更低更穩？」",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 右上區塊：KPI 定義 =====
    add_content_box(
        slide, Inches(8.85), Inches(0.85), Inches(4.2), Inches(2.1),
        "量測 KPI",
        [
            "端到端延遲 (Input→Display / Click-to-Photon)",
            "",
            "輸出指標：",
            "• P50 / P95 / P99 (ms)",
            "• 統計 + CSV（可跨版本比較）",
            "• 可選擇對齊 Perfetto / AGI trace",
            "",
            "工具可量產化：長時間自動取樣、ROI 設定",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 中間：AMD FLM vs MTK 對照表 =====
    table_title = slide.shapes.add_textbox(Inches(0.25), Inches(3.0), Inches(6), Inches(0.3))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "與 AMD Frame Latency Meter 結構相似性對照"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Microsoft JhengHei"

    table = slide.shapes.add_table(5, 3, Inches(0.25), Inches(3.28), Inches(6.5), Inches(1.55)).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.6)

    headers = ["量測模組", "AMD FLM (PC)", "MTK Latency Meter (Android)"]
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=10, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = ACCENT_BLUE

    data = [
        ["Input Trigger", "產生/捕捉輸入事件", "產生/捕捉 touch/gesture（自動化可重播）"],
        ["Frame Capture", "連續擷取畫面", "連續擷取畫面（ROI 為主）"],
        ["ROI Detector", "ROI 變化偵測（差分/門檻）", "ROI 變化偵測（差分/SAD/門檻）"],
        ["Output", "統計 + CSV", "統計 + CSV + Perfetto/AGI trace"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=9, bold=True)
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    # ===== 右側：POC 設計 =====
    add_content_box(
        slide, Inches(6.85), Inches(3.0), Inches(6.2), Inches(1.9),
        "POC 設計",
        [
            "實驗條件：",
            "A：現行方案（不量 input lag）",
            "B：現行方案 + MTK Latency Meter（外掛式/不改遊戲）",
            "C：現行方案 + MTK Latency Meter（引擎整合式/有 marker）",
            "",
            "遊戲場景：① 遊戲高刷（穩態） ② 遊戲 + 抖音複合場景",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 底部左：成功判定 =====
    add_content_box(
        slide, Inches(0.25), Inches(4.95), Inches(6.5), Inches(2.3),
        "成功判定準則",
        [
            "工具可用性：",
            "• 可長時間自動取樣",
            "• 有效 sample 率高（偵測成功率/誤判率可控）",
            "• CSV 可回歸（版本間可比）",
            "",
            "KPI 產出（在 Avg FPS 相近前提下）：",
            "• 穩定輸出 Input→Display：P50 / P95 / P99 (ms)",
            "• A/B/C 下看到可重複差異（方向一致、幅度可信）",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 底部右：技術限制說明 =====
    add_content_box(
        slide, Inches(6.85), Inches(4.95), Inches(6.2), Inches(2.3),
        "技術限制與定位說明",
        [
            "Android 的 capture/權限/overlay 路徑與 PC 不同，兩者不等價",
            "",
            "但在「用 ROI 變化偵測把端到端延遲工具化」方法論上具可比性",
            "",
            "工程焦點：",
            "• 把限制說清楚，讓測試設計能穩定重複",
            "• 不追求任何場景都能量，而是讓可量的場景可回歸",
            "",
            "MTK Latency Meter 價值：不替代 FPS 指標，而是補齊「體感」KPI",
        ],
        title_color=RGBColor(100, 100, 100)
    )

    prs.save('D:/pptx/觸控延遲量測Latency_Meter.pptx')
    print("已生成：觸控延遲量測Latency_Meter.pptx")

if __name__ == "__main__":
    create_pptx()
