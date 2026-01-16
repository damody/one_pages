# -*- coding: utf-8 -*-
"""
生成 Anti-Lag SDK 降低觸控延遲 一頁式投影片
根據 Phase 5 的內容產生 PPTX
修正版：正確繪製「由上到下詳細內部流程」的前後對比圖
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import re

# =============================================================================
# 顏色定義 (MTK 風格)
# =============================================================================

BG_COLOR = RGBColor(255, 249, 230)  # 米色背景
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)

ACCENT_BLUE = RGBColor(70, 130, 180)      # 技術背景、成功判定
ACCENT_ORANGE = RGBColor(230, 126, 34)    # 問題、POC 設計
ACCENT_GREEN = RGBColor(39, 174, 96)      # 效益、解決方案
ACCENT_PURPLE = RGBColor(142, 68, 173)    # 架構、技術細節
ACCENT_RED = RGBColor(192, 0, 0)          # 風險、警告

COLOR_RED = RGBColor(244, 67, 54)
COLOR_GREEN = RGBColor(76, 175, 80)
COLOR_BLUE = RGBColor(33, 150, 243)
COLOR_GRAY_BG = RGBColor(245, 245, 245)
COLOR_TEXT = RGBColor(51, 51, 51)
COLOR_PINK = RGBColor(233, 30, 99)

FONT_NAME = "Microsoft JhengHei"

# =============================================================================
# 流程圖佈局配置與自適應邏輯
# =============================================================================

FLOW_LAYOUT_CONFIG = {
    "max_horizontal_nodes": 6,        # 超過此數量強制縱向
    "min_node_width": 0.8,            # 最小節點寬度（英吋）
    "min_gap": 0.08,                  # 最小節點間距（英吋）
    "default_gap": 0.12,              # 預設節點間距（英吋）
    "char_width_avg": 0.09,           # 平均字元寬度估算（英吋，8pt）
    "padding_horizontal": 0.2,        # 節點內水平留白（英吋）
}

def calculate_text_width(text):
    """計算文字寬度（中文=2，英文=1）"""
    width = 0
    for char in text:
        if ord(char) > 0x4E00 and ord(char) < 0x9FFF:  # CJK 範圍
            width += 2
        else:
            width += 1
    return width

def estimate_node_min_width(node):
    """估算單個節點的最小寬度"""
    if isinstance(node, dict):
        title = node.get("title", "")
        desc = node.get("desc", "")
        time_label = node.get("time", "")
    else:
        title = str(node)
        desc = ""
        time_label = ""

    max_text = max([title, desc, time_label], key=lambda x: calculate_text_width(x))
    text_width_chars = calculate_text_width(max_text)
    estimated_width = text_width_chars * FLOW_LAYOUT_CONFIG["char_width_avg"] + FLOW_LAYOUT_CONFIG["padding_horizontal"]

    return max(estimated_width, FLOW_LAYOUT_CONFIG["min_node_width"])

def should_use_vertical_flow(nodes, width, min_node_width=0.8):
    """判斷是否應該使用縱向流程圖"""
    node_count = len(nodes)

    # 條件 1：節點數量過多
    if node_count > FLOW_LAYOUT_CONFIG["max_horizontal_nodes"]:
        return (True, f"節點數量 {node_count} > {FLOW_LAYOUT_CONFIG['max_horizontal_nodes']}")

    # 條件 2：計算動態調整後的節點寬度
    min_gap = FLOW_LAYOUT_CONFIG["min_gap"]
    calculated_node_width = (width - min_gap * (node_count - 1)) / node_count

    if calculated_node_width < min_node_width:
        return (True, f"計算節點寬度 {calculated_node_width:.2f} < 最小寬度 {min_node_width}")

    # 條件 3：檢查文字長度
    max_required_width = max(estimate_node_min_width(node) for node in nodes)
    if max_required_width > calculated_node_width:
        return (True, f"文字所需寬度 {max_required_width:.2f} > 節點寬度 {calculated_node_width:.2f}")

    return (False, "橫向佈局可行")

# =============================================================================
# 元素追蹤與排版審查
# =============================================================================

slide_elements = {}
current_slide_index = 0


def reset_element_tracker():
    global slide_elements, current_slide_index
    slide_elements = {}
    current_slide_index = 0


def set_current_slide(index):
    global current_slide_index
    current_slide_index = index
    if index not in slide_elements:
        slide_elements[index] = []


def track_element(name, left, top, width, height, element_type="generic"):
    global slide_elements, current_slide_index
    if current_slide_index not in slide_elements:
        slide_elements[current_slide_index] = []
    slide_elements[current_slide_index].append({
        "name": name,
        "type": element_type,
        "left": left,
        "top": top,
        "right": left + width,
        "bottom": top + height
    })


def boxes_overlap(box1, box2):
    return not (
        box1["right"] <= box2["left"] or
        box1["left"] >= box2["right"] or
        box1["bottom"] <= box2["top"] or
        box1["top"] >= box2["bottom"]
    )


def check_overlaps(slide_index):
    if slide_index not in slide_elements:
        return []
    elements = slide_elements[slide_index]
    overlaps = []
    non_bg_elements = [e for e in elements if e["type"] != "background"]
    for i in range(len(non_bg_elements)):
        for j in range(i + 1, len(non_bg_elements)):
            box1 = non_bg_elements[i]
            box2 = non_bg_elements[j]
            if boxes_overlap(box1, box2):
                overlaps.append({
                    "element_a": box1,
                    "element_b": box2
                })
    return overlaps


def layout_review(max_rounds=2):
    print(f"\n{'='*50}")
    print(f"開始排版審查（最多 {max_rounds} 輪）")
    print(f"{'='*50}")
    total_overlaps = 0
    for slide_idx in slide_elements:
        overlaps = check_overlaps(slide_idx)
        element_count = len([e for e in slide_elements[slide_idx] if e["type"] != "background"])
        print(f"\n投影片 {slide_idx + 1}:")
        print(f"  檢查元素數量：{element_count}")
        print(f"  偵測到重疊：{len(overlaps)} 處")
        if overlaps:
            total_overlaps += len(overlaps)
            for idx, overlap in enumerate(overlaps, 1):
                ea = overlap["element_a"]
                eb = overlap["element_b"]
                print(f"\n  重疊 {idx}:")
                print(f"    「{ea['name']}」與「{eb['name']}」重疊")
    print(f"\n{'='*50}")
    if total_overlaps == 0:
        print("排版審查通過！無重疊")
        return {"passed": True, "total_overlaps": 0}
    else:
        print(f"排版審查完成：發現 {total_overlaps} 處重疊")
        return {"passed": False, "total_overlaps": total_overlaps}


# =============================================================================
# 基本元件函數
# =============================================================================

def add_background(slide, prs, color=BG_COLOR):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    return background


def add_main_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(
        Inches(0.25), Inches(0.1), Inches(12.5), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.52), Inches(12.8), Inches(0.28)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.font.name = FONT_NAME


def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.03),
        Inches(width - 0.16), Inches(0.28)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT_NAME

    content_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.28),
        Inches(width - 0.16), Inches(height - 0.32)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_NAME
        p.space_after = Pt(1)

    return shape


def add_section_title(slide, left, top, width, text, color=ACCENT_BLUE):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(0.3)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    return box


def set_cell_text(cell, text, font_size=9, bold=True, color=DARK_GRAY):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = FONT_NAME


def add_table(slide, left, top, width, height, headers, data, header_color=ACCENT_BLUE):
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=9, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = header_color

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=8, bold=True)
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    return table


# =============================================================================
# 新增：詳細版前後對比圖（由上到下的內部流程）
# =============================================================================

def draw_before_after_with_vertical_flow(slide, left, top, width, height,
                                          before_title, before_nodes, before_summary,
                                          after_title, after_nodes, after_summary,
                                          arrow_label="導入 SDK", effect_label="效果：延遲降低 5%+"):
    """
    繪製帶有垂直內部流程的前後對比圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        before_title: 左側標題
        before_nodes: 左側節點列表 [{
            "title": "節點標題",
            "desc": "說明文字",
            "color": 顏色（可選）,
            "is_problem": True/False（是否標示為問題，紅色虛線框）
        }, ...]
        before_summary: 左側底部總結文字
        after_title: 右側標題
        after_nodes: 右側節點列表
        after_summary: 右側底部總結文字
        arrow_label: 中間箭頭上方文字
        effect_label: 中間箭頭下方文字
    """
    # 計算佈局
    box_width = (width - 0.6) / 2  # 中間留 0.6 吋給箭頭
    summary_height = 0.35
    content_height = height - summary_height - 0.4  # 減去標題和底部

    # ===== 左側區塊（改善前）=====
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    # 左側標題
    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.28)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    # 左側垂直流程節點
    node_count = len(before_nodes)
    node_gap = 0.06
    node_height = (content_height - node_gap * (node_count - 1)) / node_count
    node_width = box_width - 0.2

    for i, node in enumerate(before_nodes):
        node_top = top + 0.38 + i * (node_height + node_gap)
        node_left = left + 0.1

        color = node.get("color", COLOR_BLUE)
        is_problem = node.get("is_problem", False)

        # 節點矩形
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(node_left), Inches(node_top),
            Inches(node_width), Inches(node_height)
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = color
        node_shape.line.fill.background()

        # 如果是問題節點，加上紅色虛線框
        if is_problem:
            problem_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(node_left - 0.03), Inches(node_top - 0.03),
                Inches(node_width + 0.06), Inches(node_height + 0.06)
            )
            problem_box.fill.background()
            problem_box.line.color.rgb = COLOR_RED
            problem_box.line.width = Pt(2)
            problem_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # 節點文字
        tf = node_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = node.get("title", "")
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if node.get("desc"):
            p2 = tf.add_paragraph()
            p2.text = node.get("desc")
            p2.font.size = Pt(6)
            p2.font.color.rgb = RGBColor(255, 255, 200) if is_problem else WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        # 節點之間的箭頭
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(node_left + node_width/2 - 0.06), Inches(node_top + node_height),
                Inches(0.12), Inches(node_gap)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()

    # 左側底部總結
    before_summary_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + height - summary_height - 0.05),
        Inches(box_width - 0.2), Inches(summary_height)
    )
    tf = before_summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = before_summary
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    # ===== 右側區塊（改善後）=====
    after_left = left + box_width + 0.6
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(after_left), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    # 右側標題
    after_title_box = slide.shapes.add_textbox(
        Inches(after_left + 0.1), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.28)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    # 右側垂直流程節點
    after_node_count = len(after_nodes)
    after_node_height = (content_height - node_gap * (after_node_count - 1)) / after_node_count

    for i, node in enumerate(after_nodes):
        node_top = top + 0.38 + i * (after_node_height + node_gap)
        node_left = after_left + 0.1

        color = node.get("color", COLOR_GREEN)

        # 節點矩形
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(node_left), Inches(node_top),
            Inches(node_width), Inches(after_node_height)
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = color
        node_shape.line.fill.background()

        # 節點文字
        tf = node_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = node.get("title", "")
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if node.get("desc"):
            p2 = tf.add_paragraph()
            p2.text = node.get("desc")
            p2.font.size = Pt(6)
            p2.font.color.rgb = WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        # 節點之間的箭頭
        if i < after_node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(node_left + node_width/2 - 0.06), Inches(node_top + after_node_height),
                Inches(0.12), Inches(node_gap)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()

    # 右側底部總結
    after_summary_box = slide.shapes.add_textbox(
        Inches(after_left + 0.1), Inches(top + height - summary_height - 0.05),
        Inches(box_width - 0.2), Inches(summary_height)
    )
    tf = after_summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = after_summary
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    # ===== 中間箭頭 =====
    arrow_x = left + box_width + 0.15
    arrow_y = top + height / 2 - 0.15

    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(arrow_x), Inches(arrow_y),
        Inches(0.3), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()

    # 箭頭上方文字
    arrow_label_box = slide.shapes.add_textbox(
        Inches(arrow_x - 0.05), Inches(arrow_y - 0.25),
        Inches(0.4), Inches(0.22)
    )
    tf = arrow_label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = arrow_label
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 箭頭下方文字
    effect_label_box = slide.shapes.add_textbox(
        Inches(arrow_x - 0.1), Inches(arrow_y + 0.35),
        Inches(0.5), Inches(0.25)
    )
    tf = effect_label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = effect_label
    p.font.size = Pt(6)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


# =============================================================================
# 其他圖表繪製函數
# =============================================================================

def draw_flow(slide, left, top, width, height, nodes):
    """繪製橫向流程圖（動態調整版本）"""
    node_count = len(nodes)

    # 動態計算最佳 gap
    total_min_width = sum(estimate_node_min_width(node) for node in nodes)
    available_space = width - total_min_width

    if available_space < 0:
        gap = FLOW_LAYOUT_CONFIG["min_gap"]
    else:
        optimal_gap = available_space / (node_count - 1) if node_count > 1 else 0
        gap = min(optimal_gap, FLOW_LAYOUT_CONFIG["default_gap"])

    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            time_label = node.get("time", "")
            color = node.get("color", COLOR_BLUE)
        else:
            title = str(node)
            desc = ""
            time_label = ""
            color = COLOR_BLUE

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(7)
            p2.font.color.rgb = WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if time_label:
            p3 = tf.add_paragraph()
            p3.text = time_label
            p3.font.size = Pt(6)
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(255, 255, 200)
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER

        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.01), Inches(top + height/2 - 0.06),
                Inches(0.1), Inches(0.12)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()


def draw_flow_vertical(slide, left, top, width, height, nodes):
    """縱向流程圖（後備方案）"""
    node_count = len(nodes)
    node_gap = 0.06
    node_height = (height - node_gap * (node_count - 1)) / node_count
    node_width = width - 0.2

    for i, node in enumerate(nodes):
        node_top = top + i * (node_height + node_gap)
        node_left = left + 0.1

        # 解析節點內容
        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            time_label = node.get("time", "")
            color = node.get("color", COLOR_BLUE)
            highlight = node.get("highlight", False)
        else:
            title = str(node)
            desc = ""
            time_label = ""
            color = COLOR_BLUE
            highlight = False

        # 繪製節點
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(node_left), Inches(node_top),
            Inches(node_width), Inches(node_height)
        )
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = color
        node_shape.line.fill.background()

        # 高亮標記
        if highlight:
            highlight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(node_left - 0.03), Inches(node_top - 0.03),
                Inches(node_width + 0.06), Inches(node_height + 0.06)
            )
            highlight_box.fill.background()
            highlight_box.line.color.rgb = COLOR_RED
            highlight_box.line.width = Pt(2)
            highlight_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # 文字內容
        tf = node_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(6)
            p2.font.color.rgb = WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if time_label:
            p3 = tf.add_paragraph()
            p3.text = time_label
            p3.font.size = Pt(6)
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(255, 255, 200)
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER

        # 下箭頭
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(node_left + node_width/2 - 0.06), Inches(node_top + node_height),
                Inches(0.12), Inches(node_gap)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()


def draw_flow_adaptive(slide, left, top, width, height, nodes, arrow_labels=None):
    """自適應流程圖：優先橫向，必要時自動切換縱向"""
    should_vertical, reason = should_use_vertical_flow(nodes, width)

    if should_vertical:
        print(f"[自動切換縱向] {reason}")
        return draw_flow_vertical(slide, left, top, width, height, nodes)
    else:
        return draw_flow(slide, left, top, width, height, nodes)


def draw_architecture(slide, left, top, width, height, layers):
    """繪製分層架構圖"""
    layer_count = len(layers)
    layer_height = (height - 0.08 * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + 0.08)
        color = layer.get("color", COLOR_BLUE)
        is_highlight = layer.get("highlight", False)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        if is_highlight:
            box.line.color.rgb = ACCENT_ORANGE
            box.line.width = Pt(3)

        name_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(y + layer_height/2 - 0.12),
            Inches(1.0), Inches(0.24)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer.get("name", "")
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        components = layer.get("components", [])
        if components:
            comp_width = (width - 1.3) / len(components)
            for j, comp in enumerate(components):
                cx = left + 1.2 + j * comp_width
                comp_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(y + 0.06),
                    Inches(comp_width - 0.08), Inches(layer_height - 0.12)
                )
                comp_box.fill.solid()
                comp_box.fill.fore_color.rgb = WHITE
                comp_box.line.fill.background()

                tf = comp_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = comp
                p.font.size = Pt(7)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER


def draw_platform_compare(slide, left, top, width, height, platform_a, platform_b):
    """繪製上下平台對比圖"""
    platform_height = (height - 0.15) / 2

    platforms = [platform_a, platform_b]
    colors = [COLOR_BLUE, ACCENT_GREEN]

    for i, platform in enumerate(platforms):
        y = top + i * (platform_height + 0.15)
        color = platform.get("color", colors[i])
        name = platform.get("name", f"平台 {i+1}")
        title = platform.get("title", "")
        items = platform.get("items", [])
        summary = platform.get("summary", "")

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(platform_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_GRAY_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)

        name_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.1), Inches(y + 0.08),
            Inches(0.9), Inches(0.24)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = color
        name_box.line.fill.background()

        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(
            Inches(left + 1.1), Inches(y + 0.08),
            Inches(width - 1.3), Inches(0.24)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = FONT_NAME

        if items:
            content_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(y + 0.38),
                Inches(width - 0.2), Inches(platform_height - 0.6)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(7)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME

        if summary:
            summary_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(y + platform_height - 0.22),
                Inches(width - 0.2), Inches(0.18)
            )
            tf = summary_box.text_frame
            p = tf.paragraphs[0]
            p.text = summary
            p.font.size = Pt(7)
            p.font.bold = True
            p.font.color.rgb = color
            p.font.name = FONT_NAME


# =============================================================================
# 術語卡片函數
# =============================================================================

def draw_glossary_card(slide, left, top, width, height, term, desc):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_BG
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(0.5)

    term_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.05),
        Inches(width - 0.16), Inches(0.26)
    )
    tf = term_box.text_frame
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.31),
        Inches(width - 0.16), Inches(height - 0.38)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME


def draw_glossary_page(slide, title, terms):
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    margin = 0.25
    gap = 0.1
    cols = 4
    rows = 4
    card_width = (13.333 - margin * 2 - gap * (cols - 1)) / cols
    card_height = (7.5 - 0.5 - gap * (rows - 1)) / rows

    for i, term_data in enumerate(terms[:16]):
        row = i // cols
        col = i % cols
        x = margin + col * (card_width + gap)
        y = 0.5 + row * (card_height + gap)

        draw_glossary_card(
            slide, x, y, card_width, card_height,
            term_data.get("term", ""),
            term_data.get("desc", "")
        )


# =============================================================================
# 主程式
# =============================================================================

def create_pptx():
    reset_element_tracker()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # 第 1 頁：主投影片
    # =========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(0)
    add_background(slide1, prs)

    # 主標題
    add_main_title(
        slide1,
        "導入 Anti-Lag SDK 可望降低手遊觸控延遲 5% 以上",
        "靈感來源：AMD Radeon Anti-Lag 2 SDK，已在 PC 平台實證降低 37% 延遲 | 目標：透過 SDK 實現遊戲引擎與 SoC 的深度同步"
    )
    track_element("主標題", 0.25, 0.1, 12.5, 0.7, "text")

    # ===== 主圖：詳細的前後對比圖（由上到下的內部流程）=====
    add_section_title(slide1, 0.25, 0.85, 6.5, "延遲鏈路對比", ACCENT_BLUE)

    draw_before_after_with_vertical_flow(
        slide=slide1,
        left=0.25, top=1.12, width=6.5, height=3.5,
        before_title="改善前：傳統延遲路徑",
        before_nodes=[
            {"title": "觸控事件", "desc": "手指觸碰螢幕 T=0", "color": COLOR_PINK},
            {"title": "輸入採樣", "desc": "立即採樣（可能過時）", "color": COLOR_BLUE, "is_problem": True},
            {"title": "CPU 處理", "desc": "CPU 太快產生排隊", "color": COLOR_BLUE, "is_problem": True},
            {"title": "Frame Queue", "desc": "2-3 幀排隊等 GPU", "color": COLOR_RED, "is_problem": True},
            {"title": "GPU 渲染", "desc": "渲染排隊中的舊幀", "color": RGBColor(97, 97, 97)},
            {"title": "螢幕顯示", "desc": "顯示過時畫面", "color": RGBColor(97, 97, 97)}
        ],
        before_summary="總延遲 = 輸入 + Queue + 渲染",
        after_title="改善後：Anti-Lag SDK 同步",
        after_nodes=[
            {"title": "觸控事件", "desc": "手指觸碰螢幕 T=0", "color": COLOR_PINK},
            {"title": "SDK 同步點", "desc": "等 GPU 準備好再採樣", "color": COLOR_GREEN},
            {"title": "輸入採樣", "desc": "最佳時機採樣", "color": COLOR_GREEN},
            {"title": "CPU+GPU 同步", "desc": "節奏對齊，無排隊", "color": COLOR_GREEN},
            {"title": "螢幕顯示", "desc": "顯示最新畫面", "color": COLOR_GREEN}
        ],
        after_summary="總延遲 = 輸入 + 渲染（無 Queue）",
        arrow_label="導入SDK",
        effect_label="延遲↓5%+"
    )
    track_element("延遲對比圖", 0.25, 0.85, 6.5, 3.77, "diagram")

    # 右上：已驗證的成功要素
    add_content_box(
        slide1, 6.85, 0.85, 6.2, 1.75,
        "已驗證的成功要素（PC 平台）",
        [
            "AMD Anti-Lag 2 在 PC 上已驗證有效：",
            "• 解決 Frame Queue 排隊（CPU 不無限超前 GPU）",
            "• 精準同步點設計（Present() 前插入同步）",
            "• 輸入採樣對齊至該幀渲染前最後一刻",
            "實測效果：CS2 延遲降低約 37%",
        ],
        title_color=ACCENT_BLUE
    )
    track_element("已驗證的成功要素", 6.85, 0.85, 6.2, 1.75, "card")

    # 右中：現況問題
    add_content_box(
        slide1, 6.85, 2.7, 6.2, 1.95,
        "現況與問題（天璣平台）",
        [
            "天璣 9500 + Android 的延遲瓶頸：",
            "• BufferQueue 機制天生 1-2 幀緩衝延遲",
            "• FPSGO 強於保幀率，難感知引擎內部時序",
            "• 觸控採樣與 VSync 同步比 PC 滑鼠更複雜",
            "",
            "最有效解法：導入 SDK 級別握手機制",
        ],
        title_color=ACCENT_ORANGE
    )
    track_element("現況與問題", 6.85, 2.7, 6.2, 1.95, "card")

    # 底部左：POC 設計
    add_content_box(
        slide1, 0.25, 4.75, 4.2, 2.45,
        "POC 設計",
        [
            "實驗條件：",
            "• A 組：現行方案（僅 FPSGO）",
            "• B 組：導入 MTK Anti-Lag SDK",
            "",
            "量測：FLM/ftrace/LTR 監控 Click-to-Photon",
            "場景：遊戲高刷穩態 / 遊戲+抖音複合",
            "資源：2 人 × 4 週",
        ],
        title_color=ACCENT_ORANGE
    )
    track_element("POC 設計", 0.25, 4.75, 4.2, 2.45, "card")

    # 底部中：成功判定準則
    add_content_box(
        slide1, 4.55, 4.75, 4.2, 2.45,
        "成功判定準則",
        [
            "1. 系統延遲下降 ≥ 5%",
            "   （PC 達 37%，手機保守設定）",
            "",
            "2. 1% Low 提升",
            "   Frame Time 抖動收斂",
            "",
            "3. 功耗差異 ±5% 內視為持平",
        ],
        title_color=ACCENT_BLUE
    )
    track_element("成功判定準則", 4.55, 4.75, 4.2, 2.45, "card")

    # 底部右上：風險
    add_content_box(
        slide1, 8.85, 4.75, 4.2, 1.15,
        "風險與退場機制",
        [
            "• < 3% 改善 → 終止專案",
            "• 3-5% → 視投入產出比決定",
            "• ≥ 5% → POC 成功",
        ],
        title_color=ACCENT_RED
    )
    track_element("風險與退場機制", 8.85, 4.75, 4.2, 1.15, "card")

    # 底部右下：行動建議
    action_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.85), Inches(6.0),
        Inches(4.2), Inches(1.2)
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    action_box.line.color.rgb = ACCENT_GREEN
    action_box.line.width = Pt(2)

    action_content = slide1.shapes.add_textbox(Inches(8.95), Inches(6.05), Inches(4.0), Inches(1.1))
    tf = action_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "建議行動"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = FONT_NAME

    p = tf.add_paragraph()
    p.text = "核准 POC 驗證（2人×4週）"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    p = tf.add_paragraph()
    p.text = "遊戲廠商已初步接觸，反應正面"
    p.font.size = Pt(8)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME
    track_element("行動建議", 8.85, 6.0, 4.2, 1.2, "card")

    # =========================================================================
    # 第 2 頁：附錄 - Touch-to-Display 全鏈路時間軸
    # =========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(1)
    add_background(slide2, prs)

    title_box = slide2.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 A：Touch-to-Display 全鏈路時間軸"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    add_section_title(slide2, 0.3, 0.6, 12.7, "從觸控到螢幕顯示的完整延遲鏈路")
    draw_flow(
        slide=slide2,
        left=0.3, top=0.95, width=12.7, height=1.2,
        nodes=[
            {"title": "觸控感應", "desc": "IRQ中斷", "time": "T=0", "color": COLOR_PINK},
            {"title": "事件分發", "desc": "InputDispatcher", "time": "~2ms", "color": RGBColor(97, 97, 97)},
            {"title": "遊戲引擎", "desc": "GameThread", "time": "~4ms", "color": COLOR_BLUE},
            {"title": "GPU渲染", "desc": "RenderThread", "time": "~16ms", "color": RGBColor(97, 97, 97)},
            {"title": "系統合成", "desc": "SurfaceFlinger", "time": "~28ms", "color": RGBColor(97, 97, 97)},
            {"title": "螢幕顯示", "desc": "MIPI-DSI", "time": "~33ms", "color": COLOR_GREEN}
        ]
    )
    track_element("時間軸流程", 0.3, 0.6, 12.7, 1.55, "diagram")

    add_content_box(
        slide2, 0.3, 2.3, 4.2, 2.2,
        "輸入階段（IRQ → InputDispatcher）",
        [
            "1. 觸控 IC 偵測電容變化",
            "2. 發出中斷請求 (IRQ)",
            "3. input_report_abs() 寫入座標",
            "4. InputReader 透過 EventHub 讀取",
            "5. InputDispatcher 找到焦點 App",
            "6. 透過 Unix Socket 發送事件",
        ],
        title_color=COLOR_PINK
    )

    add_content_box(
        slide2, 4.6, 2.3, 4.2, 2.2,
        "遊戲處理階段（Engine）",
        [
            "Unity: UnityMain 單執行緒處理",
            "Unreal: GameThread + RenderThread",
            "        + RHIThread 三執行緒架構",
            "",
            "引擎調用 dequeueBuffer() 申請 Buffer",
            "GPU 執行 vkQueueSubmit / glDrawArrays",
        ],
        title_color=COLOR_BLUE
    )

    add_content_box(
        slide2, 8.9, 2.3, 4.2, 2.2,
        "顯示階段（SF → Display）",
        [
            "VSync 觸發 SurfaceFlinger",
            "latchBuffer() 檢查 acquireFence",
            "HWC 決定合成策略",
            "presentDisplay() 產生 presentFence",
            "Display Driver 透過 MIPI-DSI",
            "傳輸信號驅動螢幕發光",
        ],
        title_color=COLOR_GREEN
    )

    add_section_title(slide2, 0.3, 4.65, 12.7, "遊戲引擎架構對比")
    add_content_box(
        slide2, 0.3, 4.95, 6.3, 2.25,
        "Unity - Low Latency Pattern",
        [
            "特點：",
            "• 緩衝區深度較淺（Double Buffering）",
            "• 路徑較短",
            "• 易受主執行緒波動影響",
            "",
            "適合：對延遲敏感的競技遊戲",
        ],
        title_color=COLOR_BLUE
    )

    add_content_box(
        slide2, 6.7, 4.95, 6.3, 2.25,
        "Unreal - Throughput Pattern",
        [
            "特點：",
            "• 採用 Pipelined 設計",
            "• 能提高 GPU 利用率",
            "• presentFence 到 Touch Input 時間跨度較長",
            "",
            "適合：畫面複雜的 3A 級遊戲",
        ],
        title_color=ACCENT_PURPLE
    )

    # =========================================================================
    # 第 3 頁：附錄 - PC vs 手機平台對比
    # =========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(2)
    add_background(slide3, prs)

    title_box = slide3.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 B：PC vs 手機平台延遲治理對比"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    draw_platform_compare(
        slide3, 0.3, 0.65, 8.5, 3.5,
        platform_a={
            "name": "PC 平台",
            "title": "AMD Anti-Lag 2（已驗證有效）",
            "color": COLOR_BLUE,
            "items": [
                "輸入方式：滑鼠/鍵盤，輪詢採樣 1000Hz",
                "同步 API：Vulkan / DX12 原生支援",
                "同步點：In-engine SDK，Present() 前等待 GPU",
                "顯示路徑：直接輸出到顯示器"
            ],
            "summary": "效果：CS2 延遲降低 37%"
        },
        platform_b={
            "name": "手機平台",
            "title": "MTK Anti-Lag SDK（規劃中）",
            "color": ACCENT_GREEN,
            "items": [
                "輸入方式：觸控，需與觸控面板採樣率同步",
                "同步 API：MAGT / Performance Hint API 擴展",
                "同步點：RenderThread 提交 ↔ SurfaceFlinger 消費",
                "顯示路徑：經過 SurfaceFlinger 合成"
            ],
            "summary": "目標：延遲降低 5%+"
        }
    )

    add_content_box(
        slide3, 8.9, 0.65, 4.1, 3.5,
        "主要差異",
        [
            "輸入採樣",
            "  PC: 滑鼠輪詢",
            "  手機: 觸控需與面板同步",
            "",
            "同步介面",
            "  PC: Vulkan/DX12",
            "  手機: Performance Hint API",
            "",
            "顯示路徑",
            "  PC: 直接顯示",
            "  手機: 經 SurfaceFlinger",
            "",
            "關鍵差異：手機多了",
            "BufferQueue + SF 合成層",
        ],
        title_color=ACCENT_ORANGE
    )

    add_section_title(slide3, 0.3, 4.3, 12.7, "延遲治理對照表")
    add_table(
        slide3, 0.3, 4.6, 12.7, 1.3,
        ["比較項目", "PC 平台 (AMD Anti-Lag 2)", "手機平台 (MTK Anti-Lag SDK)"],
        [
            ["同步層級", "In-engine SDK (Vulkan/DX12)", "MAGT / Performance Hint API 擴展"],
            ["關鍵同步點", "CPU Present() 與 GPU 執行完畢", "RenderThread 提交與 SurfaceFlinger 消費"],
            ["量測工具", "FLM (Software/Hardware)", "FLM + Android Systrace (AGI)"],
        ],
        header_color=ACCENT_BLUE
    )

    add_content_box(
        slide3, 0.3, 6.05, 12.7, 1.15,
        "共同目標",
        [
            "• 消除 Frame Queue 堆積，讓 CPU 和 GPU 節奏同步",
            "• 對齊輸入採樣時機，確保玩家操作對應最新畫面",
            "• 透過 SDK 實現引擎與平台的深度協作，而非單靠驅動層調整",
        ],
        title_color=ACCENT_GREEN
    )

    # =========================================================================
    # 第 4 頁：附錄 - 系統架構圖
    # =========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(3)
    add_background(slide4, prs)

    title_box = slide4.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 C：系統架構圖 - Anti-Lag SDK 介入位置"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME

    draw_architecture(
        slide4, 0.3, 0.65, 8.5, 4.5,
        layers=[
            {
                "name": "應用層",
                "color": COLOR_BLUE,
                "components": ["遊戲 App (Unity/Unreal)", "GameThread", "RenderThread"],
                "highlight": False
            },
            {
                "name": "SDK 層",
                "color": ACCENT_ORANGE,
                "components": ["MTK Anti-Lag SDK ★", "MAGT", "Perf Hint API"],
                "highlight": True
            },
            {
                "name": "框架層",
                "color": ACCENT_GREEN,
                "components": ["SurfaceFlinger", "BufferQueue", "HWC"],
                "highlight": False
            },
            {
                "name": "硬體層",
                "color": ACCENT_PURPLE,
                "components": ["GPU", "Display Controller", "MIPI-DSI"],
                "highlight": False
            }
        ]
    )

    add_content_box(
        slide4, 8.9, 0.65, 4.1, 2.1,
        "★ SDK 介入位置",
        [
            "MTK Anti-Lag SDK 負責：",
            "",
            "• 接收引擎同步訊號",
            "• 協調輸入採樣時機",
            "• 控制 Frame Queue 深度",
            "• 與 SurfaceFlinger 握手",
        ],
        title_color=ACCENT_ORANGE
    )

    add_content_box(
        slide4, 8.9, 2.85, 4.1, 2.3,
        "資料流向",
        [
            "↓ queueBuffer + acquireFence",
            "  App → BufferQueue",
            "",
            "↓ presentDisplay + presentFence",
            "  SF → HWC",
            "",
            "↓ releaseFence",
            "  允許 App 重新 dequeueBuffer",
        ],
        title_color=ACCENT_GREEN
    )

    add_section_title(slide4, 0.3, 5.3, 12.7, "SDK 實作要點")
    add_content_box(
        slide4, 0.3, 5.6, 6.3, 1.6,
        "遊戲引擎整合",
        [
            "• 遊戲需主動呼叫 SDK 設定目標幀率",
            "• 在渲染迴圈中插入同步點",
            "• 配合 RenderThread 提交時機",
            "• 需遊戲廠商配合整合（已初步接觸）",
        ],
        title_color=COLOR_BLUE
    )

    add_content_box(
        slide4, 6.7, 5.6, 6.3, 1.6,
        "平台層支援",
        [
            "• 擴展 Performance Hint API",
            "• SurfaceFlinger 感知同步訊號",
            "• 調整 BufferQueue 緩衝策略",
            "• 與 FPSGO 協同運作（不衝突）",
        ],
        title_color=ACCENT_GREEN
    )

    # =========================================================================
    # 第 5 頁：術語詞彙表
    # =========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(4)
    add_background(slide5, prs)

    glossary_terms = [
        {"term": "Click-to-Photon", "desc": "從手指點擊螢幕到畫面像素實際改變的總時間。這是衡量操作延遲的核心指標。"},
        {"term": "Frame Queue", "desc": "CPU 產生但 GPU 尚未渲染完成的畫面佇列。堆積越多，玩家操作對應的畫面越過時。"},
        {"term": "CPU", "desc": "中央處理器，負責遊戲邏輯運算。處理速度太快會導致 Frame Queue 堆積。"},
        {"term": "GPU", "desc": "圖形處理器，負責畫面渲染。Anti-Lag SDK 讓 CPU 等 GPU 準備好再送新工作。"},
        {"term": "SoC", "desc": "系統單晶片，整合 CPU、GPU 等元件。天璣 9500 是聯發科的旗艦 SoC。"},
        {"term": "BufferQueue", "desc": "Android 的畫面緩衝區機制，天生有 1-2 幀延遲，是手機平台特有的延遲來源。"},
        {"term": "SurfaceFlinger", "desc": "Android 畫面合成服務，負責將多個 App 的畫面合成為最終顯示內容。"},
        {"term": "RenderThread", "desc": "Android App 負責渲染的執行緒，SDK 的同步點就在這裡與 SurfaceFlinger 交接。"},
        {"term": "FPSGO", "desc": "聯發科的幀率穩定技術，強項是保幀率與省電，但難以主動協調延遲。"},
        {"term": "VSync", "desc": "垂直同步訊號，是螢幕刷新的節拍器。觸控採樣需與 VSync 同步。"},
        {"term": "P99 Frame Time", "desc": "99% 畫面都比這個時間快。數字越小代表最卡的時候也不會太卡。"},
        {"term": "1% Low", "desc": "最差 1% 時間內的平均幀率。比平均 FPS 更能反映卡頓程度。"},
        {"term": "MAGT", "desc": "聯發科提供給遊戲引擎的優化介面，Anti-Lag SDK 透過它與引擎溝通。"},
        {"term": "Performance Hint API", "desc": "Android 標準的效能提示介面，讓 App 告訴系統效能需求。"},
        {"term": "FLM", "desc": "Frame Latency Meter，延遲量測工具。POC 用它量測改善效果。"},
        {"term": "FPS", "desc": "每秒畫面數。POC 需在同等 FPS 下比較延遲，確保改善不是靠犧牲幀率換來。"},
    ]

    draw_glossary_page(slide5, "附錄 D：術語詞彙表", glossary_terms)

    # =========================================================================
    # 排版審查
    # =========================================================================
    review_result = layout_review(max_rounds=2)

    # 儲存
    output_path = "D:/pptx/skills/output/phase6/one_page.pptx"
    prs.save(output_path)
    print(f"\n完整投影片已生成：{output_path}")
    print("包含 5 頁投影片：")
    print("  1. 主投影片 - Anti-Lag SDK 降低觸控延遲 POC 提案（含詳細流程對比圖）")
    print("  2. 附錄 A - Touch-to-Display 全鏈路時間軸")
    print("  3. 附錄 B - PC vs 手機平台對比")
    print("  4. 附錄 C - 系統架構圖")
    print("  5. 附錄 D - 術語詞彙表")

    if not review_result["passed"]:
        print(f"\n⚠️ 警告：排版審查發現 {review_result['total_overlaps']} 處重疊")
    else:
        print("\n✓ 排版審查通過！無元素重疊")

    return prs


if __name__ == "__main__":
    create_pptx()
