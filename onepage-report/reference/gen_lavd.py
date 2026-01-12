# -*- coding: utf-8 -*-
"""
生成 LAVD 排程 一頁式投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 顏色定義
BG_COLOR = RGBColor(255, 249, 230)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_ORANGE = RGBColor(230, 126, 34)
ACCENT_GREEN = RGBColor(39, 174, 96)
ACCENT_PURPLE = RGBColor(142, 68, 173)

def set_cell_text(cell, text, font_size=10, bold=True, color=DARK_GRAY):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = "Microsoft JhengHei"

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft JhengHei"

    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.16), height - Inches(0.32))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(1)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()

    # ===== 標題區 =====
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FPSGO + LAVD 互補方案 - 解決 1% Low 掉幀問題"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft JhengHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "靈感來源：Steam Deck scx_lavd | 目標：以 tail latency 為 KPI，同算力下提升 1% low"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 左上：歷史回顧 =====
    add_content_box(
        slide, Inches(0.25), Inches(0.85), Inches(4.25), Inches(2.0),
        "歷史回顧",
        [
            "MTK 已有 FPSGO：",
            "• 有效做幀率/功耗導向的資源治理",
            "",
            "但在複雜遊戲情境仍有限制：",
            "• 關鍵線程過多、互相喚醒/等待關係複雜",
            "• 僅靠 FPSGO 難做 task graph 級精細排程",
            "• 1% low / P99 frame time 仍可能不穩",
            "",
            "Steam Deck 已驗證：LAVD 主打 tail latency",
        ],
        title_color=ACCENT_PURPLE
    )

    # ===== 中上：問題痛點 =====
    add_content_box(
        slide, Inches(4.55), Inches(0.85), Inches(4.25), Inches(2.0),
        "問題痛點",
        [
            "遊戲一幀由多個 thread 串成 pipeline：",
            "Game → Render → RHI → Driver → OS",
            "",
            "掉幀通常來自：",
            "• 少數關鍵幀被延後 (tail latency)",
            "• 不是平均算力不足",
            "",
            "FPSGO 限制：",
            "• 偏向幀率/負載/功耗調控迴路",
            "• 只靠 boost/調頻容易變成粗粒度",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 右上：解法主張 =====
    add_content_box(
        slide, Inches(8.85), Inches(0.85), Inches(4.2), Inches(2.0),
        "解法主張：分工互補",
        [
            "FPSGO：回答「資源要給多少」",
            "• Macro：頻率 / target FPS / 功耗策略",
            "",
            "LAVD：回答「資源先給誰」",
            "• Micro：排程順序 / 關鍵喚醒延遲",
            "",
            "功能開關：",
            "透過 MAGT 下 hint 開關此 Feature",
            "僅在有效益的遊戲場景啟用",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 中間：分工對照表 =====
    table_title = slide.shapes.add_textbox(Inches(0.25), Inches(2.95), Inches(6), Inches(0.3))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "FPSGO vs LAVD 分工對照"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Microsoft JhengHei"

    table = slide.shapes.add_table(4, 3, Inches(0.25), Inches(3.2), Inches(6.5), Inches(1.25)).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.6)

    headers = ["面向", "FPSGO", "LAVD"]
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=10, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = ACCENT_BLUE

    data = [
        ["角色", "資源治理 (Macro)", "排程控制 (Micro)"],
        ["核心問題", "資源要給多少？", "資源先給誰？"],
        ["手段", "頻率/target FPS/功耗策略", "排程順序/關鍵喚醒延遲"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=9, bold=True)
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    # ===== 右中：預期增量 =====
    add_content_box(
        slide, Inches(6.85), Inches(2.95), Inches(6.2), Inches(1.55),
        "預期增量",
        [
            "在同樣算力（相近平均 util/freq）下：",
            "",
            "• Frame critical path thread 更少被插隊",
            "• 降低 wakeup → running latency",
            "• P99 frame time 收斂、1% low 提升",
            "• 減少「為了救尾端而暴力拉頻」的需求",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 底部左：POC 設計 =====
    add_content_box(
        slide, Inches(0.25), Inches(4.55), Inches(6.5), Inches(2.7),
        "POC 設計",
        [
            "實驗條件：",
            "• A：現行方案 (FPSGO/Loom)",
            "• B：現行方案 (Loom) + LAVD (scx_lavd)",
            "• C：現行方案 (FPSGO) + LAVD (scx_lavd)",
            "",
            "遊戲場景：",
            "① 遊戲高刷（穩態）",
            "② 遊戲 + 抖音複合場景（背景負載干擾）",
            "",
            "觀測重點：",
            "• 關鍵 thread 的 wakeup latency 分布",
            "• P99 frame time 變化 / 功耗波動趨勢",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 底部右：成功判定 =====
    add_content_box(
        slide, Inches(6.85), Inches(4.55), Inches(6.2), Inches(2.7),
        "成功判定準則",
        [
            "在 Avg FPS 相近（或 target FPS 相同）前提下：",
            "",
            "1. 1% Low 明顯提升",
            "   Frame Time 穩定度改善",
            "",
            "2. 關鍵 Thread 延遲下降",
            "   wakeup latency / runnable delay 顯著下降",
            "",
            "3. 功耗波動不惡化",
            "   理想是下降；至少不因救 delay 而急拉",
            "",
            "定位：LAVD 是 FPSGO 的互補層（非取代）",
            "以量測數據證明：同畫質同 target FPS 下，1% low 提升",
        ],
        title_color=ACCENT_BLUE
    )

    prs.save('D:/pptx/LAVD排程解決1%low掉幀問題.pptx')
    print("已生成：LAVD排程解決1%low掉幀問題.pptx")

if __name__ == "__main__":
    create_pptx()
