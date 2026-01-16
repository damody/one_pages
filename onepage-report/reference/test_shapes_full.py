"""
完整測試 PPTX Shapes 繪圖 + 文字排版
模擬真實的一頁投影片佈局
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime, timedelta

# === 顏色定義 ===
COLOR_RED = RGBColor(244, 67, 54)
COLOR_GREEN = RGBColor(76, 175, 80)
COLOR_BLUE = RGBColor(33, 150, 243)
COLOR_ORANGE = RGBColor(255, 152, 0)
COLOR_PURPLE = RGBColor(156, 39, 176)
COLOR_GRAY_BG = RGBColor(245, 245, 245)
COLOR_GRAY_DARK = RGBColor(97, 97, 97)
COLOR_TEXT = RGBColor(51, 51, 51)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(0, 121, 107)  # Teal

FONT_NAME = "Microsoft JhengHei"


# ============================================================================
# 元素追蹤與排版審查
# ============================================================================

# 全域元素追蹤清單（按投影片分組）
slide_elements = {}  # {slide_index: [{"name": ..., "type": ..., "left": ..., "top": ..., "right": ..., "bottom": ...}, ...]}
current_slide_index = 0


def reset_element_tracker():
    """重置元素追蹤清單"""
    global slide_elements, current_slide_index
    slide_elements = {}
    current_slide_index = 0


def set_current_slide(index):
    """設定當前投影片索引"""
    global current_slide_index
    current_slide_index = index
    if index not in slide_elements:
        slide_elements[index] = []


def track_element(name, left, top, width, height, element_type="generic"):
    """
    追蹤元素位置

    Args:
        name: 元素名稱（用於識別）
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        element_type: 元素類型 ("text", "diagram", "card", "background")
    """
    global slide_elements, current_slide_index

    if current_slide_index not in slide_elements:
        slide_elements[current_slide_index] = []

    slide_elements[current_slide_index].append({
        "name": name,
        "type": element_type,
        "left": left,
        "top": top,
        "right": left + width,
        "bottom": top + height
    })


def boxes_overlap(box1, box2):
    """
    檢查兩個元素是否重疊

    Args:
        box1, box2: {"left": float, "top": float, "right": float, "bottom": float}

    Returns:
        bool: 是否重疊
    """
    return not (
        box1["right"] <= box2["left"] or   # box1 在 box2 左邊
        box1["left"] >= box2["right"] or   # box1 在 box2 右邊
        box1["bottom"] <= box2["top"] or   # box1 在 box2 上方
        box1["top"] >= box2["bottom"]      # box1 在 box2 下方
    )


def calculate_overlap_area(box1, box2):
    """計算兩個元素的重疊面積"""
    if not boxes_overlap(box1, box2):
        return 0.0

    overlap_left = max(box1["left"], box2["left"])
    overlap_right = min(box1["right"], box2["right"])
    overlap_top = max(box1["top"], box2["top"])
    overlap_bottom = min(box1["bottom"], box2["bottom"])

    return (overlap_right - overlap_left) * (overlap_bottom - overlap_top)


def check_overlaps(slide_index):
    """
    檢查指定投影片上的元素是否有重疊

    Returns:
        list: 重疊的元素對 [{"element_a": ..., "element_b": ..., "overlap_area": ...}, ...]
    """
    if slide_index not in slide_elements:
        return []

    elements = slide_elements[slide_index]
    overlaps = []

    # 過濾掉背景元素
    non_bg_elements = [e for e in elements if e["type"] != "background"]

    for i in range(len(non_bg_elements)):
        for j in range(i + 1, len(non_bg_elements)):
            box1 = non_bg_elements[i]
            box2 = non_bg_elements[j]

            if boxes_overlap(box1, box2):
                overlap_area = calculate_overlap_area(box1, box2)
                overlaps.append({
                    "element_a": box1,
                    "element_b": box2,
                    "overlap_area": overlap_area
                })

    return overlaps


def layout_review(max_rounds=2):
    """
    執行排版審查

    Args:
        max_rounds: 最大審查輪數

    Returns:
        dict: 審查結果 {"passed": bool, "rounds": int, "total_overlaps": int, "details": [...]}
    """
    print(f"\n{'='*50}")
    print(f"開始排版審查（最多 {max_rounds} 輪）")
    print(f"{'='*50}")

    all_details = []
    total_overlaps = 0

    for slide_idx in slide_elements:
        overlaps = check_overlaps(slide_idx)
        element_count = len([e for e in slide_elements[slide_idx] if e["type"] != "background"])

        print(f"\n投影片 {slide_idx + 1}:")
        print(f"  檢查元素數量：{element_count}")
        print(f"  偵測到重疊：{len(overlaps)} 處")

        if overlaps:
            total_overlaps += len(overlaps)
            for idx, overlap in enumerate(overlaps, 1):
                ea = overlap["element_a"]
                eb = overlap["element_b"]
                area = overlap["overlap_area"]

                detail = {
                    "slide": slide_idx + 1,
                    "element_a": ea["name"],
                    "element_b": eb["name"],
                    "overlap_area": area
                }
                all_details.append(detail)

                print(f"\n  重疊 {idx}:")
                print(f"    「{ea['name']}」與「{eb['name']}」重疊")
                print(f"    - {ea['name']}: ({ea['left']:.2f}, {ea['top']:.2f}) - ({ea['right']:.2f}, {ea['bottom']:.2f})")
                print(f"    - {eb['name']}: ({eb['left']:.2f}, {eb['top']:.2f}) - ({eb['right']:.2f}, {eb['bottom']:.2f})")
                print(f"    - 重疊面積：{area:.3f} 平方吋")

    print(f"\n{'='*50}")
    if total_overlaps == 0:
        print("排版審查通過！無重疊")
        print(f"{'='*50}\n")
        return {"passed": True, "rounds": 1, "total_overlaps": 0, "details": []}
    else:
        print(f"排版審查完成：發現 {total_overlaps} 處重疊")
        print("請手動調整元素位置以消除重疊")
        print(f"{'='*50}\n")
        return {"passed": False, "rounds": 1, "total_overlaps": total_overlaps, "details": all_details}


# ============================================================================
# 圖表繪製函數
# ============================================================================

def draw_before_after(slide, left, top, width, height, before_title, before_items, after_title, after_items):
    """繪製前後對比圖"""
    box_width = (width - 0.4) / 2

    # 左側區塊
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    before_content = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(box_width - 0.2), Inches(height - 0.5)
    )
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(before_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

    # 右側區塊
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.4), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    after_title_box = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    after_content = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.4),
        Inches(box_width - 0.2), Inches(height - 0.5)
    )
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(after_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

    # 中間箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.1), Inches(top + height/2 - 0.12),
        Inches(0.2), Inches(0.24)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()


def draw_flow(slide, left, top, width, height, nodes):
    """繪製橫向流程圖"""
    node_count = len(nodes)
    gap = 0.12
    arrow_width = 0.1
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            color = node.get("color", COLOR_BLUE)
        else:
            title = str(node)
            desc = ""
            color = COLOR_BLUE

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(8)
            p2.font.color.rgb = COLOR_WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.01), Inches(top + height/2 - 0.06),
                Inches(arrow_width), Inches(0.12)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()


def draw_architecture(slide, left, top, width, height, layers):
    """
    繪製分層架構圖

    Args:
        layers: [{"name": "...", "color": ..., "components": ["...", ...]}, ...]
    """
    layer_count = len(layers)
    layer_height = (height - 0.1 * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + 0.1)
        color = layer.get("color", COLOR_BLUE)

        # 層背景
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # 層名稱（左側）
        name_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(y + layer_height/2 - 0.15),
            Inches(1.2), Inches(0.3)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer.get("name", "")
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        # 組件（右側橫排）
        components = layer.get("components", [])
        if components:
            comp_width = (width - 1.5) / len(components)
            for j, comp in enumerate(components):
                cx = left + 1.4 + j * comp_width
                comp_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(y + 0.1),
                    Inches(comp_width - 0.1), Inches(layer_height - 0.2)
                )
                comp_box.fill.solid()
                comp_box.fill.fore_color.rgb = COLOR_WHITE
                comp_box.line.fill.background()

                tf = comp_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = comp
                p.font.size = Pt(8)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER


def draw_metric_cards(slide, left, top, width, height, metrics):
    """
    繪製指標卡片（數字 + 說明）

    Args:
        metrics: [{"value": "81%", "label": "延遲降低", "color": ...}, ...]
    """
    card_count = len(metrics)
    card_width = (width - 0.15 * (card_count - 1)) / card_count

    for i, metric in enumerate(metrics):
        x = left + i * (card_width + 0.15)
        color = metric.get("color", COLOR_BLUE)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(card_width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        # 數值（大字）
        value_box = slide.shapes.add_textbox(
            Inches(x), Inches(top + 0.1),
            Inches(card_width), Inches(height * 0.5)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get("value", "")
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 標籤（小字）
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(top + height * 0.55),
            Inches(card_width), Inches(height * 0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get("label", "")
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER


def draw_comparison_table(slide, left, top, width, height, headers, rows):
    """
    繪製對比表格

    Args:
        headers: ["項目", "PC", "手機"]
        rows: [["延遲", "20ms", "50ms"], ...]
    """
    col_count = len(headers)
    row_count = len(rows) + 1  # +1 for header
    col_width = width / col_count
    row_height = height / row_count

    for r in range(row_count):
        for c in range(col_count):
            x = left + c * col_width
            y = top + r * row_height

            # 儲存格背景
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(col_width), Inches(row_height)
            )

            if r == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BLUE
                text_color = COLOR_WHITE
                is_bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_GRAY_BG if r % 2 == 1 else COLOR_WHITE
                text_color = COLOR_TEXT
                is_bold = False

            cell.line.color.rgb = RGBColor(200, 200, 200)
            cell.line.width = Pt(0.5)

            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(y + 0.05),
                Inches(col_width - 0.1), Inches(row_height - 0.1)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            if r == 0:
                p.text = headers[c]
            else:
                p.text = rows[r-1][c]

            p.font.size = Pt(9)
            p.font.bold = is_bold
            p.font.color.rgb = text_color
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER


def draw_icon_list(slide, left, top, width, item_height, items):
    """
    繪製帶圖標的列表

    Args:
        items: [{"icon": "check/cross/warn", "text": "...", "color": ...}, ...]
    """
    for i, item in enumerate(items):
        y = top + i * item_height
        icon_type = item.get("icon", "check")
        color = item.get("color", COLOR_GREEN if icon_type == "check" else COLOR_RED)

        # 圖標
        icon_size = 0.2
        if icon_type == "check":
            # 綠色圓圈 + 勾
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_GREEN
            circle.line.fill.background()
        elif icon_type == "cross":
            # 紅色圓圈 + 叉
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_RED
            circle.line.fill.background()
        elif icon_type == "warn":
            # 黃色三角形
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            triangle.fill.solid()
            triangle.fill.fore_color.rgb = COLOR_ORANGE
            triangle.line.fill.background()
        else:  # 藍色圓點
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left), Inches(y + (item_height - icon_size) / 2),
                Inches(icon_size), Inches(icon_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_BLUE
            circle.line.fill.background()

        # 文字
        text_box = slide.shapes.add_textbox(
            Inches(left + icon_size + 0.1), Inches(y),
            Inches(width - icon_size - 0.1), Inches(item_height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("text", "")
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME


# ============================================================================
# 原生圖表繪製函數（使用 python-pptx XL_CHART_TYPE）
# ============================================================================

def draw_line_chart(slide, left, top, width, height, title, categories, series,
                    show_legend=True, show_data_labels=False, smooth=False):
    """
    繪製折線圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        categories: X 軸類別 ["Week 1", "Week 2", ...]
        series: 資料序列 [{"name": "延遲", "values": [80, 65, 45, 30], "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_data_labels: 是否顯示資料標籤
        smooth: 是否平滑曲線
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series:
        chart_data.add_series(s["name"], s["values"])

    chart_type = XL_CHART_TYPE.LINE_MARKERS_STACKED if not smooth else XL_CHART_TYPE.LINE_MARKERS

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    # 設定標題
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    # 設定圖例
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # 設定顏色
    for i, s in enumerate(series):
        if "color" in s and i < len(chart.series):
            chart.series[i].format.line.color.rgb = s["color"]

    return graphic_frame


def draw_bar_chart(slide, left, top, width, height, title, categories, series,
                   show_legend=True, show_data_labels=False, horizontal=False):
    """
    繪製長條圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        categories: X 軸類別 ["項目A", "項目B", ...]
        series: 資料序列 [{"name": "數據", "values": [10, 20, 30], "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_data_labels: 是否顯示資料標籤
        horizontal: 是否水平顯示
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series:
        chart_data.add_series(s["name"], s["values"])

    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    # 設定標題
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    # 設定圖例
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # 設定顏色
    for i, s in enumerate(series):
        if "color" in s and i < len(chart.series):
            for point in chart.series[i].points:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = s["color"]

    return graphic_frame


def draw_pie_chart(slide, left, top, width, height, title, data,
                   show_legend=True, show_percentage=True):
    """
    繪製圓餅圖（原生圖表）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        data: 資料 [{"name": "項目", "value": 45, "color": ...}, ...]
        show_legend: 是否顯示圖例
        show_percentage: 是否顯示百分比
    """
    chart_data = CategoryChartData()
    chart_data.categories = [d["name"] for d in data]
    chart_data.add_series("Values", [d["value"] for d in data])

    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    # 設定標題
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME

    # 設定圖例
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    # 設定顏色
    if len(chart.series) > 0:
        for i, d in enumerate(data):
            if "color" in d and i < len(chart.series[0].points):
                chart.series[0].points[i].format.fill.solid()
                chart.series[0].points[i].format.fill.fore_color.rgb = d["color"]

    return graphic_frame


# ============================================================================
# Shapes 組合圖表（甘特圖、矩陣圖、增強版架構圖）
# ============================================================================

def draw_gantt_chart(slide, left, top, width, height, title, tasks, milestones=None,
                     time_unit="week", show_today_line=True, show_progress=True):
    """
    繪製甘特圖（使用 Shapes 組合）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        tasks: 任務列表 [{"name": "任務", "start": "2025-01-06", "end": "2025-01-19",
                         "progress": 100, "color": COLOR_BLUE}, ...]
        milestones: 里程碑 [{"name": "里程碑", "date": "2025-02-28", "color": COLOR_RED}, ...]
        time_unit: 時間單位 "day" | "week" | "month"
        show_today_line: 是否顯示今日線
        show_progress: 是否顯示進度條
    """
    # 解析日期
    def parse_date(date_str):
        return datetime.strptime(date_str, "%Y-%m-%d")

    # 計算時間範圍
    all_dates = []
    for t in tasks:
        all_dates.append(parse_date(t["start"]))
        all_dates.append(parse_date(t["end"]))
    if milestones:
        for m in milestones:
            all_dates.append(parse_date(m["date"]))

    min_date = min(all_dates)
    max_date = max(all_dates)
    total_days = (max_date - min_date).days + 1

    # 佈局參數
    title_height = 0.35
    header_height = 0.3
    name_col_width = 1.8
    chart_left = left + name_col_width
    chart_width = width - name_col_width
    task_count = len(tasks)
    task_height = (height - title_height - header_height) / task_count if task_count > 0 else 0.5

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(title_height)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 時間軸標題區
    header_top = top + title_height

    # 計算時間刻度
    if time_unit == "week":
        # 每週一個刻度
        current = min_date
        ticks = []
        while current <= max_date:
            ticks.append(current)
            current += timedelta(days=7)
        if ticks[-1] < max_date:
            ticks.append(max_date)
    else:
        # 每天或每月（簡化為等分 5-8 個刻度）
        tick_count = min(8, total_days)
        ticks = [min_date + timedelta(days=i * total_days // tick_count) for i in range(tick_count + 1)]

    # 繪製時間刻度
    for i, tick_date in enumerate(ticks[:-1]):
        tick_x = chart_left + (tick_date - min_date).days / total_days * chart_width
        next_tick_x = chart_left + (ticks[i + 1] - min_date).days / total_days * chart_width if i + 1 < len(ticks) else chart_left + chart_width

        # 刻度標籤
        label_box = slide.shapes.add_textbox(
            Inches(tick_x), Inches(header_top),
            Inches(next_tick_x - tick_x), Inches(header_height)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = tick_date.strftime("%m/%d")
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_GRAY_DARK
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 垂直分隔線
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(tick_x), Inches(header_top + header_height),
            Inches(0.01), Inches(height - title_height - header_height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(220, 220, 220)
        line.line.fill.background()

    # 繪製任務
    task_area_top = header_top + header_height
    for i, task in enumerate(tasks):
        task_top = task_area_top + i * task_height
        task_start = parse_date(task["start"])
        task_end = parse_date(task["end"])
        color = task.get("color", COLOR_BLUE)
        progress = task.get("progress", 0)

        # 任務名稱
        name_box = slide.shapes.add_textbox(
            Inches(left), Inches(task_top),
            Inches(name_col_width - 0.1), Inches(task_height)
        )
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = task["name"]
        p.font.size = Pt(8)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

        # 計算任務條位置
        start_offset = (task_start - min_date).days / total_days
        end_offset = (task_end - min_date).days / total_days
        bar_left = chart_left + start_offset * chart_width
        bar_width = (end_offset - start_offset) * chart_width
        bar_top = task_top + task_height * 0.2
        bar_height = task_height * 0.6

        # 任務背景條
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_left), Inches(bar_top),
            Inches(bar_width), Inches(bar_height)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(230, 230, 230)
        bg_bar.line.fill.background()

        # 進度條
        if show_progress and progress > 0:
            progress_width = bar_width * progress / 100
            progress_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(bar_top),
                Inches(progress_width), Inches(bar_height)
            )
            progress_bar.fill.solid()
            progress_bar.fill.fore_color.rgb = color
            progress_bar.line.fill.background()

    # 繪製里程碑
    if milestones:
        for m in milestones:
            m_date = parse_date(m["date"])
            m_offset = (m_date - min_date).days / total_days
            m_x = chart_left + m_offset * chart_width
            m_color = m.get("color", COLOR_RED)

            # 菱形
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(m_x - 0.1), Inches(task_area_top - 0.05),
                Inches(0.2), Inches(0.2)
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = m_color
            diamond.line.fill.background()

    # 今日線
    if show_today_line:
        today = datetime.now()
        if min_date <= today <= max_date:
            today_offset = (today - min_date).days / total_days
            today_x = chart_left + today_offset * chart_width
            today_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(today_x), Inches(header_top),
                Inches(0.02), Inches(height - title_height)
            )
            today_line.fill.solid()
            today_line.fill.fore_color.rgb = COLOR_RED
            today_line.line.fill.background()


def draw_architecture_enhanced(slide, left, top, width, height, layers,
                                connections=None, annotations=None):
    """
    繪製增強版系統架構圖

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        layers: 分層定義 [{"name": "層名", "color": ..., "components": [
            {"name": "組件", "highlight": False, "highlight_label": ""}, ...]}, ...]
        connections: 連接線 [{"from": "組件A", "to": "組件B", "label": "...", "style": "arrow"}, ...]
        annotations: 註解 [{"target": "組件", "text": "說明", "position": "right"}, ...]
    """
    component_positions = {}  # 記錄每個組件的位置

    layer_count = len(layers)
    layer_gap = 0.15
    layer_height = (height - layer_gap * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + layer_gap)
        color = layer.get("color", COLOR_BLUE)
        name = layer.get("name", "")

        # 層背景
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        layer_box.fill.solid()
        # 使用淺灰色背景搭配有色邊框
        layer_box.fill.fore_color.rgb = COLOR_GRAY_BG
        layer_box.line.color.rgb = color
        layer_box.line.width = Pt(2)

        # 層名稱（左側標籤）
        name_width = 1.2
        name_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.05), Inches(y + 0.05),
            Inches(name_width - 0.1), Inches(layer_height - 0.1)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = color
        name_box.line.fill.background()

        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 組件
        components = layer.get("components", [])
        if components:
            comp_area_width = width - name_width - 0.2
            comp_width = comp_area_width / len(components) if len(components) > 0 else 1
            comp_height = layer_height - 0.2

            for j, comp in enumerate(components):
                comp_name = comp["name"] if isinstance(comp, dict) else comp
                highlight = comp.get("highlight", False) if isinstance(comp, dict) else False
                highlight_label = comp.get("highlight_label", "") if isinstance(comp, dict) else ""

                cx = left + name_width + 0.1 + j * comp_width
                cy = y + 0.1

                # 組件方塊
                comp_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(cy),
                    Inches(comp_width - 0.1), Inches(comp_height)
                )
                comp_box.fill.solid()
                comp_box.fill.fore_color.rgb = COLOR_WHITE
                comp_box.line.color.rgb = color
                comp_box.line.width = Pt(1)

                # 高亮處理
                if highlight:
                    comp_box.line.color.rgb = COLOR_ORANGE
                    comp_box.line.width = Pt(3)

                    # 高亮標籤
                    if highlight_label:
                        label_box = slide.shapes.add_textbox(
                            Inches(cx), Inches(cy - 0.18),
                            Inches(comp_width - 0.1), Inches(0.18)
                        )
                        tf = label_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = highlight_label
                        p.font.size = Pt(6)
                        p.font.bold = True
                        p.font.color.rgb = COLOR_ORANGE
                        p.font.name = FONT_NAME
                        p.alignment = PP_ALIGN.CENTER

                # 組件名稱
                tf = comp_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = comp_name
                p.font.size = Pt(8)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER

                # 記錄位置
                component_positions[comp_name] = {
                    "left": cx,
                    "top": cy,
                    "width": comp_width - 0.1,
                    "height": comp_height,
                    "center_x": cx + (comp_width - 0.1) / 2,
                    "center_y": cy + comp_height / 2,
                    "bottom": cy + comp_height
                }

    # 繪製連接線
    if connections:
        for conn in connections:
            from_pos = component_positions.get(conn["from"])
            to_pos = component_positions.get(conn["to"])

            if from_pos and to_pos:
                # 計算連線的起點和終點
                start_x = from_pos["center_x"]
                start_y = from_pos["bottom"]
                end_x = to_pos["center_x"]
                end_y = to_pos["top"]

                # 使用連接器形狀
                style = conn.get("style", "arrow")
                label = conn.get("label", "")

                # 垂直線
                if abs(start_x - end_x) < 0.1:
                    line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(start_x - 0.01), Inches(start_y),
                        Inches(0.02), Inches(end_y - start_y)
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = COLOR_GRAY_DARK
                    line.line.fill.background()
                else:
                    # 斜線（用小箭頭表示）
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        Inches(min(start_x, end_x)), Inches(start_y),
                        Inches(abs(end_x - start_x) + 0.1), Inches(end_y - start_y)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = COLOR_GRAY_DARK
                    arrow.line.fill.background()

                # 連線標籤
                if label:
                    label_x = (start_x + end_x) / 2 - 0.3
                    label_y = (start_y + end_y) / 2 - 0.1
                    label_box = slide.shapes.add_textbox(
                        Inches(label_x), Inches(label_y),
                        Inches(0.6), Inches(0.2)
                    )
                    tf = label_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(6)
                    p.font.color.rgb = COLOR_GRAY_DARK
                    p.font.name = FONT_NAME
                    p.alignment = PP_ALIGN.CENTER

    return component_positions


def draw_matrix_chart(slide, left, top, width, height, title, x_axis, y_axis, items, quadrant_labels=None):
    """
    繪製矩陣圖（風險評估/優先級矩陣）

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        title: 圖表標題
        x_axis: {"label": "影響程度", "values": ["低", "中", "高"]}
        y_axis: {"label": "發生機率", "values": ["低", "中", "高"]}
        items: [{"name": "風險項", "x": 2, "y": 1, "color": COLOR_ORANGE}, ...]
        quadrant_labels: 象限標籤 {"high-high": "高優先", "low-low": "低優先", ...}
    """
    title_height = 0.35
    axis_label_width = 0.4
    axis_value_height = 0.25

    chart_left = left + axis_label_width + 0.1
    chart_top = top + title_height
    chart_width = width - axis_label_width - 0.2
    chart_height = height - title_height - axis_value_height - 0.1

    x_count = len(x_axis["values"])
    y_count = len(y_axis["values"])
    cell_width = chart_width / x_count
    cell_height = chart_height / y_count

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(title_height)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # Y 軸標籤
    y_label_box = slide.shapes.add_textbox(
        Inches(left), Inches(chart_top + chart_height / 2 - 0.2),
        Inches(axis_label_width), Inches(0.4)
    )
    tf = y_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = y_axis["label"]
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # X 軸標籤
    x_label_box = slide.shapes.add_textbox(
        Inches(chart_left + chart_width / 2 - 0.4), Inches(top + height - 0.25),
        Inches(0.8), Inches(0.25)
    )
    tf = x_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = x_axis["label"]
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 繪製格子和軸值
    for yi, y_val in enumerate(reversed(y_axis["values"])):  # 從下到上
        for xi, x_val in enumerate(x_axis["values"]):
            cell_x = chart_left + xi * cell_width
            cell_y = chart_top + yi * cell_height

            # 計算背景色（越高越紅）
            risk_level = (xi + (y_count - 1 - yi)) / (x_count + y_count - 2)
            bg_r = int(232 + (255 - 232) * risk_level)
            bg_g = int(245 - (245 - 200) * risk_level)
            bg_b = int(233 - (233 - 200) * risk_level)

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cell_x), Inches(cell_y),
                Inches(cell_width), Inches(cell_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
            cell.line.color.rgb = RGBColor(200, 200, 200)
            cell.line.width = Pt(0.5)

    # Y 軸刻度值
    for yi, y_val in enumerate(reversed(y_axis["values"])):
        val_box = slide.shapes.add_textbox(
            Inches(left + axis_label_width - 0.3), Inches(chart_top + yi * cell_height + cell_height / 2 - 0.1),
            Inches(0.3), Inches(0.2)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = y_val
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.RIGHT

    # X 軸刻度值
    for xi, x_val in enumerate(x_axis["values"]):
        val_box = slide.shapes.add_textbox(
            Inches(chart_left + xi * cell_width), Inches(chart_top + chart_height),
            Inches(cell_width), Inches(axis_value_height)
        )
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = x_val
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # 繪製項目點
    for item in items:
        ix = item["x"]  # 0-indexed
        iy = item["y"]  # 0-indexed
        color = item.get("color", COLOR_ORANGE)

        # 計算位置（y 軸需要反轉）
        point_x = chart_left + (ix + 0.5) * cell_width - 0.15
        point_y = chart_top + (y_count - 1 - iy + 0.5) * cell_height - 0.15

        # 項目圓點
        point = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(point_x), Inches(point_y),
            Inches(0.3), Inches(0.3)
        )
        point.fill.solid()
        point.fill.fore_color.rgb = color
        point.line.color.rgb = COLOR_WHITE
        point.line.width = Pt(2)

        # 項目名稱
        name_box = slide.shapes.add_textbox(
            Inches(point_x - 0.2), Inches(point_y + 0.32),
            Inches(0.7), Inches(0.2)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = item["name"]
        p.font.size = Pt(6)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER


# ============================================================================
# 進階圖表繪製函數（支援詳細內容）
# ============================================================================

def draw_flow_detailed(slide, left, top, width, height, nodes, arrow_labels=None, show_highlight=True):
    """
    繪製詳細版橫向流程圖，支援箭頭上的文字標籤和高亮節點

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        nodes: 節點列表，每個元素是 {
            "title": "節點標題",
            "desc": "說明文字（可選）",
            "time": "時間標籤（可選）",
            "color": 顏色（可選，預設 COLOR_BLUE）,
            "highlight": True/False（是否高亮，用紅色虛線框）
        }
        arrow_labels: 箭頭上的文字標籤列表（長度應為 len(nodes)-1）
        show_highlight: 是否顯示高亮標記
    """
    node_count = len(nodes)
    gap = 0.18  # 節點間距（箭頭空間）
    arrow_width = 0.15
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
            time_label = node.get("time", "")
            color = node.get("color", COLOR_BLUE)
            highlight = node.get("highlight", False)
        else:
            title = str(node)
            desc = ""
            time_label = ""
            color = COLOR_BLUE
            highlight = False

        # 節點矩形
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # 高亮標記（紅色虛線框）
        if highlight and show_highlight:
            highlight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x - 0.03), Inches(top - 0.03),
                Inches(node_width + 0.06), Inches(height + 0.06)
            )
            highlight_box.fill.background()
            highlight_box.line.color.rgb = COLOR_RED
            highlight_box.line.width = Pt(2)
            highlight_box.line.dash_style = 2  # dash

        # 節點文字
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(7)
            p2.font.color.rgb = COLOR_WHITE
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.CENTER

        if time_label:
            p3 = tf.add_paragraph()
            p3.text = time_label
            p3.font.size = Pt(7)
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(255, 255, 200)  # 淺黃色
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER

        # 箭頭（最後一個不加）
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.02), Inches(top + height/2 - 0.08),
                Inches(arrow_width), Inches(0.16)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.fill.background()

            # 箭頭上的文字標籤
            if arrow_labels and i < len(arrow_labels) and arrow_labels[i]:
                label_box = slide.shapes.add_textbox(
                    Inches(x + node_width + 0.02), Inches(top - 0.18),
                    Inches(gap - 0.04), Inches(0.18)
                )
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = arrow_labels[i]
                p.font.size = Pt(6)
                p.font.color.rgb = COLOR_TEXT
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER


def draw_before_after_with_flow(slide, left, top, width, height,
                                 before_title, before_flow_nodes, before_arrow_labels,
                                 after_title, after_flow_nodes, after_arrow_labels,
                                 center_arrow_label="導入方案", bottom_table=None):
    """
    繪製帶內部流程圖的前後對比

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        before_title: 左側標題（如「改善前：畫面堆積導致延遲」）
        before_flow_nodes: 左側內部流程節點列表
        before_arrow_labels: 左側箭頭標籤列表
        after_title: 右側標題
        after_flow_nodes: 右側內部流程節點列表
        after_arrow_labels: 右側箭頭標籤列表
        center_arrow_label: 中間箭頭上的文字（如「導入 SDK」）
        bottom_table: 底部對比表格 {"headers": [...], "rows": [[...], ...]}
    """
    # 計算佈局
    table_height = 0.7 if bottom_table else 0
    main_height = height - table_height
    box_width = (width - 0.5) / 2  # 中間留 0.5 吋給箭頭
    flow_area_height = main_height - 0.5  # 減去標題高度

    # 左側區塊（改善前）
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(main_height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    # 左側標題
    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.35)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME

    # 左側內部流程圖
    if before_flow_nodes:
        draw_flow_detailed(
            slide,
            left=left + 0.1,
            top=top + 0.45,
            width=box_width - 0.2,
            height=flow_area_height - 0.1,
            nodes=before_flow_nodes,
            arrow_labels=before_arrow_labels,
            show_highlight=True
        )

    # 右側區塊（改善後）
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.5), Inches(top),
        Inches(box_width), Inches(main_height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    # 右側標題
    after_title_box = slide.shapes.add_textbox(
        Inches(left + box_width + 0.6), Inches(top + 0.08),
        Inches(box_width - 0.2), Inches(0.35)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    # 右側內部流程圖
    if after_flow_nodes:
        draw_flow_detailed(
            slide,
            left=left + box_width + 0.6,
            top=top + 0.45,
            width=box_width - 0.2,
            height=flow_area_height - 0.1,
            nodes=after_flow_nodes,
            arrow_labels=after_arrow_labels,
            show_highlight=False
        )

    # 中間箭頭
    arrow_y = top + main_height / 2 - 0.15
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.1), Inches(arrow_y),
        Inches(0.3), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()

    # 中間箭頭上的文字
    if center_arrow_label:
        center_label = slide.shapes.add_textbox(
            Inches(left + box_width + 0.05), Inches(arrow_y - 0.22),
            Inches(0.4), Inches(0.2)
        )
        tf = center_label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = center_arrow_label
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLUE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # 底部對比表格
    if bottom_table:
        headers = bottom_table.get("headers", [])
        rows = bottom_table.get("rows", [])
        if headers and rows:
            draw_comparison_table(
                slide,
                left=left, top=top + main_height + 0.05,
                width=width, height=table_height - 0.05,
                headers=headers, rows=rows
            )


def draw_platform_compare(slide, left, top, width, height,
                          platform_a, platform_b, differences=None):
    """
    繪製上下平台對比圖，每個平台內部有完整流程

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        platform_a: 上方平台 {
            "name": "PC 平台",
            "title": "AMD Anti-Lag 2 流程（已驗證有效）",
            "color": COLOR_BLUE,
            "flow_nodes": [...],
            "arrow_labels": [...],
            "summary": "效果：CS2 延遲降低 37%"
        }
        platform_b: 下方平台（同上格式）
        differences: 差異標註列表 [{"item": "輸入方式", "a": "滑鼠", "b": "觸控"}, ...]
    """
    # 計算佈局
    diff_width = 0 if not differences else 2.5  # 差異欄位寬度
    platform_width = width - diff_width - 0.1
    platform_height = (height - 0.2) / 2

    platforms = [platform_a, platform_b]
    colors = [COLOR_BLUE, COLOR_GREEN]

    for i, platform in enumerate(platforms):
        y = top + i * (platform_height + 0.2)
        color = platform.get("color", colors[i])
        name = platform.get("name", f"平台 {i+1}")
        title = platform.get("title", "")
        flow_nodes = platform.get("flow_nodes", [])
        arrow_labels = platform.get("arrow_labels", [])
        summary = platform.get("summary", "")

        # 平台區塊
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(platform_width), Inches(platform_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_GRAY_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # 平台名稱標籤（左上角）
        name_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.1), Inches(y + 0.08),
            Inches(1.0), Inches(0.28)
        )
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = color
        name_box.line.fill.background()

        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 標題
        title_box = slide.shapes.add_textbox(
            Inches(left + 1.2), Inches(y + 0.08),
            Inches(platform_width - 1.4), Inches(0.28)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = FONT_NAME

        # 內部流程圖
        if flow_nodes:
            flow_height = platform_height - 0.7 if summary else platform_height - 0.45
            draw_flow_detailed(
                slide,
                left=left + 0.1,
                top=y + 0.4,
                width=platform_width - 0.2,
                height=flow_height,
                nodes=flow_nodes,
                arrow_labels=arrow_labels,
                show_highlight=True
            )

        # 效果摘要
        if summary:
            summary_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(y + platform_height - 0.25),
                Inches(platform_width - 0.2), Inches(0.2)
            )
            tf = summary_box.text_frame
            p = tf.paragraphs[0]
            p.text = summary
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = color
            p.font.name = FONT_NAME

    # 差異標註欄（右側）
    if differences:
        diff_left = left + platform_width + 0.1
        diff_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(diff_left), Inches(top),
            Inches(diff_width), Inches(height)
        )
        diff_box.fill.solid()
        diff_box.fill.fore_color.rgb = RGBColor(255, 243, 224)  # 淺橙色背景
        diff_box.line.color.rgb = COLOR_ORANGE
        diff_box.line.width = Pt(1)

        # 差異標題
        diff_title = slide.shapes.add_textbox(
            Inches(diff_left + 0.1), Inches(top + 0.08),
            Inches(diff_width - 0.2), Inches(0.25)
        )
        tf = diff_title.text_frame
        p = tf.paragraphs[0]
        p.text = "主要差異"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_ORANGE
        p.font.name = FONT_NAME

        # 差異項目
        item_height = (height - 0.4) / len(differences)
        for j, diff in enumerate(differences):
            item_y = top + 0.35 + j * item_height
            item_box = slide.shapes.add_textbox(
                Inches(diff_left + 0.1), Inches(item_y),
                Inches(diff_width - 0.2), Inches(item_height)
            )
            tf = item_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = diff.get("item", "")
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_NAME

            p2 = tf.add_paragraph()
            p2.text = f"PC: {diff.get('a', '')}"
            p2.font.size = Pt(7)
            p2.font.color.rgb = COLOR_BLUE
            p2.font.name = FONT_NAME

            p3 = tf.add_paragraph()
            p3.text = f"手機: {diff.get('b', '')}"
            p3.font.size = Pt(7)
            p3.font.color.rgb = COLOR_GREEN
            p3.font.name = FONT_NAME


# ============================================================================
# 文字方塊輔助函數
# ============================================================================

def add_section_title(slide, left, top, width, text, color=COLOR_BLUE):
    """加入區塊標題"""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    return box


def add_bullet_list(slide, left, top, width, height, items, font_size=10):
    """加入項目列表"""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
        p.space_after = Pt(4)

    return box


def add_content_box(slide, left, top, width, height, title, content, title_color=COLOR_BLUE, bg_color=None):
    """加入帶標題的內容方塊"""
    if bg_color:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

    # 標題
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08),
        Inches(width - 0.2), Inches(0.3)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT_NAME

    # 內容
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(width - 0.2), Inches(height - 0.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    if isinstance(content, list):
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_NAME
    else:
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME


# ============================================================================
# 術語卡片函數
# ============================================================================

def draw_mini_flow(slide, left, top, width, height, nodes):
    """繪製迷你流程圖（用於術語卡片內的示意圖）"""
    node_count = len(nodes)
    gap = 0.08
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)
        color = node.get("color", COLOR_BLUE) if isinstance(node, dict) else COLOR_BLUE
        text = node.get("text", node) if isinstance(node, dict) else str(node)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # 箭頭
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.01), Inches(top + height/2 - 0.04),
                Inches(gap - 0.02), Inches(0.08)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()


def draw_mini_before_after(slide, left, top, width, height, before_text, after_text):
    """繪製迷你前後對比圖（用於術語卡片內的示意圖）"""
    box_width = (width - 0.25) / 2

    # 左側（Before）
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(1)

    tf = before_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = before_text
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_RED
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 右側（After）
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.25), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(1)

    tf = after_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = after_text
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 中間箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.05), Inches(top + height/2 - 0.06),
        Inches(0.15), Inches(0.12)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()


def draw_mini_layers(slide, left, top, width, height, layers):
    """繪製迷你分層圖（用於術語卡片內的示意圖）"""
    layer_count = len(layers)
    layer_height = (height - 0.05 * (layer_count - 1)) / layer_count

    for i, layer in enumerate(layers):
        y = top + i * (layer_height + 0.05)
        color = layer.get("color", COLOR_BLUE) if isinstance(layer, dict) else COLOR_BLUE
        text = layer.get("text", layer) if isinstance(layer, dict) else str(layer)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(layer_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER


def draw_mini_timeline(slide, left, top, width, height, stages):
    """繪製迷你時間軸（用於術語卡片內的示意圖）"""
    line_y = top + height * 0.6
    stage_count = len(stages)
    stage_width = width / stage_count

    # 水平軸線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(line_y),
        Inches(width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(100, 100, 100)
    line.line.fill.background()

    for i, stage in enumerate(stages):
        x = left + i * stage_width
        color = stage.get("color", COLOR_BLUE) if isinstance(stage, dict) else COLOR_BLUE
        text = stage.get("text", stage) if isinstance(stage, dict) else str(stage)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.02), Inches(top),
            Inches(stage_width - 0.04), Inches(height * 0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(6)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER


def draw_mini_icon(slide, left, top, width, height, icon_type, label):
    """繪製迷你圖標（用於術語卡片內的示意圖）"""
    icon_size = min(width, height * 0.6)
    icon_x = left + (width - icon_size) / 2
    icon_y = top

    if icon_type == "chart":
        # 長條圖圖標
        bar_width = icon_size / 4
        bars = [0.4, 0.7, 0.5, 0.9]
        for i, h in enumerate(bars):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(icon_x + i * bar_width), Inches(icon_y + icon_size * (1 - h)),
                Inches(bar_width - 0.02), Inches(icon_size * h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_BLUE
            bar.line.fill.background()
    elif icon_type == "gauge":
        # 儀表圖標
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_GREEN
        circle.line.fill.background()
    elif icon_type == "warning":
        # 警告圖標
        triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = COLOR_ORANGE
        triangle.line.fill.background()
    else:
        # 預設圓形
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y),
            Inches(icon_size), Inches(icon_size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_BLUE
        circle.line.fill.background()

    # 標籤
    label_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + icon_size + 0.05),
        Inches(width), Inches(height - icon_size - 0.05)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


def draw_glossary_card_with_diagram(slide, left, top, width, height, term, desc, diagram_type, diagram_params):
    """
    繪製帶示意圖的術語卡片

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        term: 術語名稱
        desc: 術語解釋
        diagram_type: 示意圖類型 (flow, before_after, layers, timeline, icon)
        diagram_params: 示意圖參數
    """
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_BG
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(1)

    # 示意圖區域（上方 50%）
    diagram_height = height * 0.45
    diagram_top = top + 0.1
    diagram_left = left + 0.15
    diagram_width = width - 0.3

    if diagram_type == "flow":
        draw_mini_flow(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("nodes", []))
    elif diagram_type == "before_after":
        draw_mini_before_after(slide, diagram_left, diagram_top, diagram_width, diagram_height,
                               diagram_params.get("before", ""), diagram_params.get("after", ""))
    elif diagram_type == "layers":
        draw_mini_layers(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("layers", []))
    elif diagram_type == "timeline":
        draw_mini_timeline(slide, diagram_left, diagram_top, diagram_width, diagram_height, diagram_params.get("stages", []))
    elif diagram_type == "icon":
        draw_mini_icon(slide, diagram_left, diagram_top, diagram_width, diagram_height,
                       diagram_params.get("icon_type", ""), diagram_params.get("label", ""))

    # 術語名稱（中間）
    term_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + diagram_height + 0.15),
        Inches(width - 0.2), Inches(0.35)
    )
    tf = term_box.text_frame
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    # 解釋文字（下方）
    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + diagram_height + 0.5),
        Inches(width - 0.2), Inches(height - diagram_height - 0.6)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME


def draw_glossary_card_text_only(slide, left, top, width, height, term, desc):
    """
    繪製純文字術語卡片

    Args:
        slide: 投影片物件
        left, top: 左上角位置（吋）
        width, height: 寬高（吋）
        term: 術語名稱
        desc: 術語解釋（簡短版）
    """
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_BG
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(0.5)

    # 術語名稱
    term_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.05),
        Inches(width - 0.16), Inches(0.28)
    )
    tf = term_box.text_frame
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    # 解釋文字
    desc_box = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 0.33),
        Inches(width - 0.16), Inches(height - 0.4)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME


def draw_glossary_page_with_diagrams(slide, title, terms):
    """
    繪製一頁 6 格有圖片的術語卡片（2 列 x 3 欄）

    Args:
        slide: 投影片物件
        title: 頁面標題
        terms: 最多 6 個術語，每個是 {
            "term": "術語名稱",
            "desc": "解釋",
            "diagram_type": "flow|before_after|layers|timeline|icon",
            "diagram_params": {...}
        }
    """
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 佈局參數
    margin = 0.3
    gap = 0.2
    cols = 3
    rows = 2
    card_width = (13.333 - margin * 2 - gap * (cols - 1)) / cols  # ≈ 4.1"
    card_height = (7.5 - 0.6 - gap * (rows - 1)) / rows  # ≈ 3.3"

    for i, term_data in enumerate(terms[:6]):
        row = i // cols
        col = i % cols
        x = margin + col * (card_width + gap)
        y = 0.6 + row * (card_height + gap)

        draw_glossary_card_with_diagram(
            slide, x, y, card_width, card_height,
            term_data.get("term", ""),
            term_data.get("desc", ""),
            term_data.get("diagram_type", "icon"),
            term_data.get("diagram_params", {})
        )


def draw_glossary_page_text_only(slide, title, terms):
    """
    繪製一頁 16 格純文字術語卡片（4 列 x 4 欄）

    Args:
        slide: 投影片物件
        title: 頁面標題
        terms: 最多 16 個術語，每個是 {
            "term": "術語名稱",
            "desc": "簡短解釋（<=50字）"
        }
    """
    # 頁面標題
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 佈局參數
    margin = 0.25
    gap = 0.1
    cols = 4
    rows = 4
    card_width = (13.333 - margin * 2 - gap * (cols - 1)) / cols  # ≈ 3.1"
    card_height = (7.5 - 0.5 - gap * (rows - 1)) / rows  # ≈ 1.7"

    for i, term_data in enumerate(terms[:16]):
        row = i // cols
        col = i % cols
        x = margin + col * (card_width + gap)
        y = 0.5 + row * (card_height + gap)

        draw_glossary_card_text_only(
            slide, x, y, card_width, card_height,
            term_data.get("term", ""),
            term_data.get("desc", "")
        )


# ============================================================================
# 主程式 - 建立完整測試投影片
# ============================================================================

if __name__ == "__main__":
    # 初始化元素追蹤
    reset_element_tracker()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========================================================================
    # 第 1 頁：模擬真實報告 - 主投影片
    # ========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(0)  # 設定當前投影片索引

    # 主標題
    title = slide1.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "導入 SDK 同步機制可降低 81% 遊戲操作延遲"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 副標題
    subtitle = slide1.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(12.7), Inches(0.3))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "借鏡 PC Anti-Lag 技術，解決手機 Frame Queue 堆積問題"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_GRAY_DARK
    p.font.name = FONT_NAME

    # 左側：問題與方案（前後對比圖）
    add_section_title(slide1, 0.3, 1.0, 6.0, "問題與解決方案")
    draw_before_after(
        slide=slide1,
        left=0.3, top=1.35, width=6.2, height=2.3,
        before_title="現況問題",
        before_items=[
            "GPU 不知道遊戲目標幀率",
            "BufferQueue 堆積 2-3 幀",
            "操作延遲達 50-80ms",
            "玩家體驗卡頓、不跟手"
        ],
        after_title="導入 SDK 後",
        after_items=[
            "SDK 通知目標幀率",
            "GPU 同步渲染節奏",
            "延遲降低至 10-15ms",
            "操作即時反饋"
        ]
    )
    track_element("問題與解決方案", 0.3, 1.0, 6.2, 2.65, "diagram")

    # 右側：關鍵指標
    add_section_title(slide1, 6.8, 1.0, 6.0, "預期效益")
    draw_metric_cards(
        slide=slide1,
        left=6.8, top=1.35, width=6.0, height=1.1,
        metrics=[
            {"value": "81%", "label": "延遲降低", "color": COLOR_GREEN},
            {"value": "0", "label": "Frame Queue", "color": COLOR_BLUE},
            {"value": "99%", "label": "幀率穩定度", "color": COLOR_ACCENT}
        ]
    )
    track_element("預期效益", 6.8, 1.0, 6.0, 1.45, "diagram")

    # 右側：驗證清單
    add_section_title(slide1, 6.8, 2.6, 6.0, "成功判定準則")
    draw_icon_list(
        slide=slide1,
        left=6.8, top=2.9, width=5.8, item_height=0.32,  # 縮小 item_height 避免重疊
        items=[
            {"icon": "check", "text": "Click-to-Photon 延遲 < 30ms"},
            {"icon": "check", "text": "Frame Queue 維持 0-1 幀"},
            {"icon": "check", "text": "功耗增加 < 5%"},
            {"icon": "warn", "text": "需驗證 10 款熱門遊戲"}
        ]
    )
    track_element("成功判定準則", 6.8, 2.6, 5.8, 1.58, "text")  # 2.6 + 0.3(標題) + 0.32*4 = 4.18

    # 中間：技術流程（調整位置避免重疊）
    add_section_title(slide1, 0.3, 4.25, 12.7, "技術方案流程")
    draw_flow(
        slide=slide1,
        left=0.3, top=4.55, width=12.7, height=0.75,  # 往下移並稍微縮小高度
        nodes=[
            {"title": "遊戲", "desc": "設定目標幀率", "color": COLOR_BLUE},
            {"title": "SDK", "desc": "傳遞同步訊號", "color": COLOR_GREEN},
            {"title": "SurfaceFlinger", "desc": "調整 Buffer", "color": COLOR_BLUE},
            {"title": "GPU Driver", "desc": "同步渲染節奏", "color": COLOR_GREEN},
            {"title": "顯示", "desc": "降低延遲", "color": COLOR_ACCENT}
        ]
    )
    track_element("技術方案流程", 0.3, 4.25, 12.7, 1.05, "diagram")

    # 下方左側：POC 設計（調整位置）
    add_section_title(slide1, 0.3, 5.45, 6.0, "POC 實驗設計")
    add_content_box(
        slide=slide1,
        left=0.3, top=5.75, width=6.2, height=1.45,  # 稍微縮小高度
        title="A/B 測試方案",
        content=[
            "A 組：現行機制（Baseline）",
            "B 組：導入 SDK 同步（Experimental）",
            "測試遊戲：王者榮耀、原神、PUBG 等 10 款",
            "量測指標：延遲、幀率、功耗、溫度"
        ],
        title_color=COLOR_BLUE,
        bg_color=COLOR_GRAY_BG
    )
    track_element("POC 實驗設計", 0.3, 5.45, 6.2, 1.75, "text")

    # 下方右側：行動建議（調整位置）
    add_section_title(slide1, 6.8, 5.45, 6.0, "建議行動")
    action_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(5.75),
        Inches(6.0), Inches(1.45)  # 稍微縮小高度
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    action_box.line.color.rgb = COLOR_GREEN
    action_box.line.width = Pt(2)

    action_content = slide1.shapes.add_textbox(Inches(6.95), Inches(5.85), Inches(5.7), Inches(1.25))
    tf = action_content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "核准 POC 開發"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.font.name = FONT_NAME

    items = [
        "第 1-2 週：SDK 原型開發",
        "第 3-4 週：整合測試與調優",
        "第 5-6 週：10 款遊戲驗證",
        "預期產出：完整評估報告與上線建議"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME
    track_element("建議行動", 6.8, 5.45, 6.0, 1.75, "text")

    # ========================================================================
    # 第 2 頁：附錄 - 架構圖 + 對比表
    # ========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(1)  # 設定當前投影片索引

    title = slide2.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 A：系統架構與平台對比"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 左側：架構圖
    add_section_title(slide2, 0.3, 0.6, 6.0, "Android 圖形系統架構")
    draw_architecture(
        slide=slide2,
        left=0.3, top=0.95, width=6.2, height=4.0,
        layers=[
            {"name": "應用層", "color": COLOR_BLUE, "components": ["遊戲 App", "SDK"]},
            {"name": "框架層", "color": COLOR_GREEN, "components": ["SurfaceFlinger", "BufferQueue"]},
            {"name": "HAL", "color": COLOR_ORANGE, "components": ["GPU Driver", "Display HAL"]},
            {"name": "硬體", "color": COLOR_PURPLE, "components": ["GPU", "Display"]}
        ]
    )

    # 左下：架構說明
    add_content_box(
        slide=slide2,
        left=0.3, top=5.1, width=6.2, height=2.1,
        title="SDK 插入位置",
        content=[
            "SDK 位於應用層，遊戲直接呼叫",
            "透過 Binder IPC 與框架層通訊",
            "SurfaceFlinger 調整 BufferQueue 行為",
            "GPU Driver 根據訊號同步渲染節奏"
        ],
        title_color=COLOR_BLUE,
        bg_color=COLOR_GRAY_BG
    )

    # 右側：對比表
    add_section_title(slide2, 6.8, 0.6, 6.0, "PC vs 手機平台對比")
    draw_comparison_table(
        slide=slide2,
        left=6.8, top=0.95, width=6.0, height=2.7,
        headers=["項目", "PC 平台", "手機平台"],
        rows=[
            ["同步機制", "Anti-Lag SDK", "待開發"],
            ["Frame Queue", "0-1 幀", "2-3 幀"],
            ["延遲", "10-20ms", "50-80ms"],
            ["功耗控制", "不重要", "關鍵指標"],
            ["SDK 介面", "成熟", "需設計"]
        ]
    )

    # 右中：重點說明
    add_content_box(
        slide=slide2,
        left=6.8, top=3.8, width=6.0, height=1.5,
        title="技術關鍵點",
        content=[
            "PC Anti-Lag 已驗證延遲可降低 81%",
            "手機平台需透過 SurfaceFlinger 整合",
            "SDK 需提供目標幀率設定介面"
        ],
        title_color=COLOR_ACCENT,
        bg_color=COLOR_GRAY_BG
    )

    # 右下：風險與挑戰
    add_content_box(
        slide=slide2,
        left=6.8, top=5.45, width=6.0, height=1.75,
        title="風險與挑戰",
        content=[
            "遊戲需主動整合 SDK（需推廣）",
            "不同 GPU 架構可能有差異",
            "功耗與延遲的權衡需調優",
            "需與各遊戲廠商合作驗證"
        ],
        title_color=COLOR_ORANGE,
        bg_color=RGBColor(255, 243, 224)
    )

    # ========================================================================
    # 第 3 頁：附錄 - 實驗設計
    # ========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(2)  # 設定當前投影片索引

    title = slide3.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 B：POC 實驗設計"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 上方：A/B 組對比
    add_section_title(slide3, 0.3, 0.6, 12.7, "實驗組別設計")

    # A 組
    add_content_box(
        slide=slide3,
        left=0.3, top=0.95, width=6.2, height=2.0,
        title="A 組 (Baseline)",
        content=[
            "現行 BufferQueue 機制",
            "無 SDK 同步介面",
            "Frame Queue 自然堆積",
            "作為對照基準"
        ],
        title_color=COLOR_RED,
        bg_color=RGBColor(255, 235, 238)
    )

    # B 組
    add_content_box(
        slide=slide3,
        left=6.8, top=0.95, width=6.2, height=2.0,
        title="B 組 (Experimental)",
        content=[
            "導入 SDK 同步機制",
            "遊戲主動設定目標幀率",
            "GPU 同步渲染節奏",
            "驗證延遲改善效果"
        ],
        title_color=COLOR_GREEN,
        bg_color=RGBColor(232, 245, 233)
    )

    # 中間：測試流程
    add_section_title(slide3, 0.3, 3.1, 12.7, "測試流程")
    draw_flow(
        slide=slide3,
        left=0.3, top=3.45, width=12.7, height=0.9,
        nodes=[
            {"title": "環境準備", "desc": "刷入測試版本"},
            {"title": "遊戲安裝", "desc": "10 款熱門遊戲"},
            {"title": "數據採集", "desc": "延遲/幀率/功耗"},
            {"title": "A/B 對比", "desc": "統計顯著性"},
            {"title": "報告產出", "desc": "結論與建議"}
        ]
    )

    # 下方左側：測試指標
    add_section_title(slide3, 0.3, 4.5, 6.0, "量測指標")
    draw_metric_cards(
        slide=slide3,
        left=0.3, top=4.85, width=6.2, height=1.0,
        metrics=[
            {"value": "延遲", "label": "Click-to-Photon", "color": COLOR_BLUE},
            {"value": "幀率", "label": "FPS & 1% Low", "color": COLOR_GREEN}
        ]
    )
    draw_metric_cards(
        slide=slide3,
        left=0.3, top=6.0, width=6.2, height=1.0,
        metrics=[
            {"value": "功耗", "label": "mW (avg/peak)", "color": COLOR_ORANGE},
            {"value": "溫度", "label": "CPU/GPU (°C)", "color": COLOR_RED}
        ]
    )

    # 下方右側：測試遊戲清單
    add_section_title(slide3, 6.8, 4.5, 6.0, "測試遊戲清單")
    add_content_box(
        slide=slide3,
        left=6.8, top=4.85, width=6.0, height=2.15,
        title="10 款熱門遊戲",
        content=[
            "MOBA：王者榮耀、傳說對決",
            "射擊：PUBG、Free Fire、使命召喚",
            "RPG：原神、崩壞：星穹鐵道",
            "競速：極速領域、QQ 飛車",
            "休閒：Candy Crush"
        ],
        title_color=COLOR_BLUE,
        bg_color=COLOR_GRAY_BG
    )

    # ========================================================================
    # 第 4 頁：術語表 - 有圖片版（6 格）
    # ========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(3)  # 設定當前投影片索引

    terms_with_diagrams = [
        {
            "term": "Frame Queue",
            "desc": "GPU 渲染完成但尚未顯示的畫面佇列。堆積越多，延遲越高。",
            "diagram_type": "before_after",
            "diagram_params": {"before": "堆積\n2-3 幀", "after": "維持\n0-1 幀"}
        },
        {
            "term": "Click-to-Photon",
            "desc": "從按下按鈕到畫面顯示的總延遲時間，競技遊戲目標 < 30ms。",
            "diagram_type": "timeline",
            "diagram_params": {"stages": [
                {"text": "輸入", "color": COLOR_BLUE},
                {"text": "處理", "color": COLOR_BLUE},
                {"text": "渲染", "color": COLOR_GREEN},
                {"text": "顯示", "color": COLOR_ACCENT}
            ]}
        },
        {
            "term": "BufferQueue",
            "desc": "Android 管理圖形緩衝區的機制，在生產者與消費者間傳遞畫面。",
            "diagram_type": "flow",
            "diagram_params": {"nodes": [
                {"text": "遊戲", "color": COLOR_BLUE},
                {"text": "Buffer", "color": COLOR_GREEN},
                {"text": "顯示", "color": COLOR_ACCENT}
            ]}
        },
        {
            "term": "SurfaceFlinger",
            "desc": "Android 畫面合成服務，將多個應用畫面合成為最終顯示內容。",
            "diagram_type": "layers",
            "diagram_params": {"layers": [
                {"text": "App 1", "color": COLOR_BLUE},
                {"text": "App 2", "color": COLOR_GREEN},
                {"text": "合成輸出", "color": COLOR_ACCENT}
            ]}
        },
        {
            "term": "Anti-Lag SDK",
            "desc": "PC 延遲優化技術，讓遊戲與 GPU 同步渲染節奏。",
            "diagram_type": "before_after",
            "diagram_params": {"before": "不同步\n高延遲", "after": "同步\n低延遲"}
        },
        {
            "term": "1% Low FPS",
            "desc": "最差的 1% 幀率，比平均幀率更能反映卡頓程度。",
            "diagram_type": "icon",
            "diagram_params": {"icon_type": "chart", "label": "幀率分布"}
        }
    ]

    draw_glossary_page_with_diagrams(slide4, "附錄 C：術語解釋（有圖片版 - 6 格）", terms_with_diagrams)

    # ========================================================================
    # 第 5 頁：術語表 - 純文字版（16 格）
    # ========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(4)  # 設定當前投影片索引

    terms_text_only = [
        {"term": "Frame Queue", "desc": "GPU 渲染完成但尚未顯示的畫面佇列"},
        {"term": "Click-to-Photon", "desc": "按下按鈕到畫面顯示的總延遲"},
        {"term": "BufferQueue", "desc": "Android 圖形緩衝區管理機制"},
        {"term": "SurfaceFlinger", "desc": "Android 畫面合成服務"},
        {"term": "Anti-Lag SDK", "desc": "PC 平台延遲優化技術"},
        {"term": "1% Low FPS", "desc": "最差 1% 幀率，衡量卡頓程度"},
        {"term": "VSync", "desc": "垂直同步信號，避免畫面撕裂"},
        {"term": "Triple Buffering", "desc": "三重緩衝，平滑幀率波動"},
        {"term": "Latency", "desc": "延遲，從輸入到輸出的時間"},
        {"term": "FPS", "desc": "每秒幀數，衡量畫面流暢度"},
        {"term": "GPU Driver", "desc": "顯示卡驅動程式"},
        {"term": "HAL", "desc": "硬體抽象層，連接系統與硬體"},
        {"term": "Pipeline", "desc": "管線，資料處理的連續階段"},
        {"term": "Render", "desc": "渲染，將資料轉換為畫面"},
        {"term": "Buffer", "desc": "緩衝區，暫存資料的記憶體"},
        {"term": "Sync", "desc": "同步，協調多個元件的時序"}
    ]

    draw_glossary_page_text_only(slide5, "附錄 D：術語速查表（純文字版 - 16 格）", terms_text_only)

    # ========================================================================
    # 第 6 頁：新增圖表類型展示（折線圖、長條圖、圓餅圖）
    # ========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(5)

    title = slide6.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 E：原生圖表類型展示"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 折線圖
    add_section_title(slide6, 0.3, 0.6, 4.0, "折線圖 - 效能趨勢")
    draw_line_chart(
        slide=slide6,
        left=0.3, top=0.95, width=4.0, height=2.8,
        title="延遲變化趨勢",
        categories=["Week 1", "Week 2", "Week 3", "Week 4", "Week 5"],
        series=[
            {"name": "改善前 (ms)", "values": [80, 75, 78, 82, 79], "color": COLOR_RED},
            {"name": "改善後 (ms)", "values": [80, 45, 30, 20, 15], "color": COLOR_GREEN}
        ],
        show_legend=True
    )

    # 長條圖
    add_section_title(slide6, 4.6, 0.6, 4.0, "長條圖 - 遊戲對比")
    draw_bar_chart(
        slide=slide6,
        left=4.6, top=0.95, width=4.0, height=2.8,
        title="各遊戲延遲對比",
        categories=["王者榮耀", "原神", "PUBG", "Free Fire"],
        series=[
            {"name": "改善前", "values": [80, 65, 70, 55], "color": COLOR_RED},
            {"name": "改善後", "values": [25, 20, 22, 18], "color": COLOR_GREEN}
        ],
        show_legend=True
    )

    # 圓餅圖
    add_section_title(slide6, 8.9, 0.6, 4.0, "圓餅圖 - 延遲來源")
    draw_pie_chart(
        slide=slide6,
        left=8.9, top=0.95, width=4.0, height=2.8,
        title="延遲來源分布",
        data=[
            {"name": "Frame Queue", "value": 45, "color": COLOR_RED},
            {"name": "GPU 渲染", "value": 30, "color": COLOR_BLUE},
            {"name": "顯示路徑", "value": 15, "color": COLOR_GREEN},
            {"name": "其他", "value": 10, "color": COLOR_GRAY_DARK}
        ],
        show_legend=True
    )

    # 矩陣圖
    add_section_title(slide6, 0.3, 3.9, 6.0, "矩陣圖 - 風險評估")
    draw_matrix_chart(
        slide=slide6,
        left=0.3, top=4.2, width=5.5, height=3.0,
        title="POC 風險評估矩陣",
        x_axis={"label": "影響程度", "values": ["低", "中", "高"]},
        y_axis={"label": "發生機率", "values": ["低", "中", "高"]},
        items=[
            {"name": "SDK 整合", "x": 1, "y": 1, "color": COLOR_ORANGE},
            {"name": "效能", "x": 2, "y": 0, "color": COLOR_GREEN},
            {"name": "廠商配合", "x": 2, "y": 2, "color": COLOR_RED},
            {"name": "功耗", "x": 0, "y": 1, "color": COLOR_BLUE}
        ]
    )

    # 說明文字
    note_box = slide6.shapes.add_textbox(Inches(6.2), Inches(4.2), Inches(6.5), Inches(3.0))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "圖表類型說明"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    p.font.name = FONT_NAME

    notes = [
        "• 折線圖：使用 XL_CHART_TYPE.LINE_MARKERS",
        "• 長條圖：使用 XL_CHART_TYPE.COLUMN_CLUSTERED",
        "• 圓餅圖：使用 XL_CHART_TYPE.PIE",
        "• 矩陣圖：使用 Shapes 組合繪製",
        "• 所有原生圖表可在 PowerPoint 中編輯資料"
    ]
    for note in notes:
        p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = FONT_NAME

    # ========================================================================
    # 第 7 頁：甘特圖與增強版架構圖
    # ========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_current_slide(6)

    title = slide7.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.4))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "附錄 F：甘特圖與增強版架構圖"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    p.font.name = FONT_NAME

    # 甘特圖
    add_section_title(slide7, 0.3, 0.6, 12.7, "甘特圖 - POC 專案進度")
    draw_gantt_chart(
        slide=slide7,
        left=0.3, top=0.95, width=12.7, height=2.8,
        title="Anti-Lag SDK POC 開發時程",
        tasks=[
            {"name": "需求分析與設計", "start": "2025-01-06", "end": "2025-01-12", "progress": 100, "color": COLOR_BLUE},
            {"name": "SDK 原型開發", "start": "2025-01-13", "end": "2025-01-26", "progress": 80, "color": COLOR_GREEN},
            {"name": "SurfaceFlinger 整合", "start": "2025-01-20", "end": "2025-02-02", "progress": 50, "color": COLOR_ORANGE},
            {"name": "整合測試", "start": "2025-02-03", "end": "2025-02-16", "progress": 20, "color": COLOR_PURPLE},
            {"name": "遊戲驗證", "start": "2025-02-10", "end": "2025-02-28", "progress": 0, "color": COLOR_ACCENT}
        ],
        milestones=[
            {"name": "原型完成", "date": "2025-01-26", "color": COLOR_BLUE},
            {"name": "POC 完成", "date": "2025-02-28", "color": COLOR_RED}
        ],
        time_unit="week",
        show_today_line=False,
        show_progress=True
    )

    # 增強版架構圖
    add_section_title(slide7, 0.3, 3.9, 12.7, "增強版架構圖 - Anti-Lag SDK 系統架構")
    draw_architecture_enhanced(
        slide=slide7,
        left=0.3, top=4.2, width=12.7, height=3.0,
        layers=[
            {
                "name": "應用層",
                "color": COLOR_BLUE,
                "components": [
                    {"name": "遊戲 App"},
                    {"name": "Anti-Lag SDK", "highlight": True, "highlight_label": "新增"},
                    {"name": "其他 App"}
                ]
            },
            {
                "name": "框架層",
                "color": COLOR_GREEN,
                "components": [
                    {"name": "SurfaceFlinger", "highlight": True, "highlight_label": "修改"},
                    {"name": "BufferQueue"},
                    {"name": "HWComposer"}
                ]
            },
            {
                "name": "驅動層",
                "color": COLOR_ORANGE,
                "components": [
                    {"name": "GPU Driver"},
                    {"name": "Display Driver"}
                ]
            }
        ],
        connections=[
            {"from": "Anti-Lag SDK", "to": "SurfaceFlinger", "label": "sync", "style": "arrow"}
        ]
    )

    # ========================================================================
    # 排版審查
    # ========================================================================
    review_result = layout_review(max_rounds=2)

    # 儲存
    output_path = "./test_shapes_full_v5.pptx"
    prs.save(output_path)
    print(f"\n完整測試完成！已儲存至：{output_path}")
    print("包含 7 頁投影片：")
    print("  1. 主投影片 - 完整報告佈局")
    print("  2. 附錄 A - 架構圖 + 對比表")
    print("  3. 附錄 B - POC 實驗設計")
    print("  4. 附錄 C - 術語解釋（有圖片版 6 格）")
    print("  5. 附錄 D - 術語速查表（純文字版 16 格）")
    print("  6. 附錄 E - 原生圖表類型展示（折線圖、長條圖、圓餅圖、矩陣圖）")
    print("  7. 附錄 F - 甘特圖與增強版架構圖")

    if not review_result["passed"]:
        print(f"\n⚠️ 警告：排版審查發現 {review_result['total_overlaps']} 處重疊")
        print("請檢查投影片並手動調整")
    else:
        print("\n✓ 排版審查通過！無元素重疊")
