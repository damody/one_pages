# -*- coding: utf-8 -*-
"""
生成 MTK Anti-Lag SDK 一頁式投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 顏色定義
BG_COLOR = RGBColor(255, 249, 230)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_ORANGE = RGBColor(230, 126, 34)
ACCENT_GREEN = RGBColor(39, 174, 96)

def set_cell_text(cell, text, font_size=10, bold=True, color=DARK_GRAY):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = "Microsoft JhengHei"

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft JhengHei"

    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.16), height - Inches(0.32))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(1)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()

    # ===== 標題區 =====
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MTK Anti-Lag SDK - 輸入-顯示延遲優化方案"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft JhengHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "靈感來源：AMD Radeon Anti-Lag 2 SDK | 目標：透過 SDK 實現引擎與 SoC 的深度節奏同步"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 左上：PC Anti-Lag 2 成功要素 =====
    add_content_box(
        slide, Inches(0.25), Inches(0.85), Inches(4.25), Inches(2.1),
        "PC Anti-Lag 2 成功要素",
        [
            "為何 AMD Anti-Lag 2 能產生顯著效果：",
            "",
            "• 解決 CPU-Bound 造成的 Frame Queue 排隊",
            "• In-engine 同步點：Present() 前插入同步",
            "• 確保 CPU 不會無限超前 GPU",
            "• Input Sampling 對齊：滑鼠/鍵盤採樣在",
            "  該幀渲染前最後一刻",
            "",
            "實測效果：CS2 延遲降低約 37%",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 中上：Android 延遲現狀 =====
    add_content_box(
        slide, Inches(4.55), Inches(0.85), Inches(4.25), Inches(2.1),
        "Dimensity / Android 延遲瓶頸",
        [
            "以天璣 9500 + Android 遊戲為例：",
            "",
            "• BufferQueue 機制：生產者(App)與消費者",
            "  (SurfaceFlinger) 天生有 1-2 幀緩衝延遲",
            "",
            "• FPSGO 限制：強於保幀率與節能，",
            "  但難感知引擎內 Input-to-Render 間隔",
            "",
            "• SoC 調度開銷：Input Event 處理與",
            "  RenderThread 喚醒抖動影響 Total Latency",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 右上：關鍵技術點 =====
    add_content_box(
        slide, Inches(8.85), Inches(0.85), Inches(4.2), Inches(2.1),
        "技術關鍵點",
        [
            "手機平台挑戰：",
            "觸控採樣(Touch Sampling)與顯示刷新率(Vsync)",
            "的同步比 PC 滑鼠更複雜",
            "",
            "唯一路徑：",
            "導入 SDK 級別的握手(Handshake)",
            "實現微秒級延遲控制",
            "",
            "同步層級：MAGT / Performance Hint API 擴展",
            "關鍵同步點：RenderThread ↔ SurfaceFlinger",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 中間：對照表 =====
    table_title = slide.shapes.add_textbox(Inches(0.25), Inches(3.0), Inches(6), Inches(0.3))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "PC vs Mobile 延遲治理對照"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Microsoft JhengHei"

    table = slide.shapes.add_table(4, 3, Inches(0.25), Inches(3.28), Inches(6.5), Inches(1.25)).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.6)

    headers = ["延遲階段", "PC (AMD Anti-Lag 2)", "Dimensity (MTK Proposed)"]
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=10, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = ACCENT_BLUE

    data = [
        ["同步層級", "In-engine SDK (Vulkan/DX12)", "MAGT / Performance Hint API 擴展"],
        ["關鍵同步點", "CPU Present() ↔ GPU 完成", "RenderThread 提交 ↔ SurfaceFlinger 消費"],
        ["量測工具", "FLM (Software/Hardware)", "FLM + Android Systrace (AGI)"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=9, bold=True)
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    # ===== 右中：預期效益 =====
    add_content_box(
        slide, Inches(6.85), Inches(3.0), Inches(6.2), Inches(1.6),
        "預期效益",
        [
            "若能實作類似 Anti-Lag 2 的 SDK 協作機制：",
            "",
            "• 消除無效排隊：GPU 忙碌時通知引擎推遲輸入採樣",
            "• 確保玩家點擊的是「螢幕上最新的一幀」",
            "• 收斂 P99 Frame Time（長尾延遲與 Frame Queue 堆疊相關）",
            "• 核心競爭力：鳴潮、和平菁英等高頻互動遊戲",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 底部左：POC 設計 =====
    add_content_box(
        slide, Inches(0.25), Inches(4.65), Inches(6.5), Inches(2.6),
        "POC 設計",
        [
            "實驗條件：",
            "• A 組 (Baseline)：現行方案（僅 FPSGO，無延遲同步）",
            "• B 組 (Experimental)：導入 MTK Anti-Lag SDK（引擎與平台同步）",
            "",
            "量測工具：",
            "• FLM / ftrace / LTR 進行 End-to-End Latency 量測",
            "• 監控 Click-to-Photon（點擊到螢幕像素變化總時間）",
            "",
            "遊戲場景：① 遊戲高刷（穩態） ② 遊戲 + 抖音複合場景",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 底部右：成功判定 =====
    add_content_box(
        slide, Inches(6.85), Inches(4.65), Inches(6.2), Inches(2.6),
        "成功判定準則",
        [
            "1. 系統延遲下降",
            "   全鏈路延遲在同等 FPS 下降低 5% 以上",
            "",
            "2. 1% Low 提升",
            "   Frame Time 抖動幅度明顯收斂",
            "",
            "3. 功耗行為",
            "   不因縮短延遲而造成溫升超標",
            "   維持功耗波動與 A 組持平",
            "",
            "工程結論：SDK 實現深度節奏同步，顯著降低「點擊到反應」總延遲",
        ],
        title_color=ACCENT_BLUE
    )

    prs.save('D:/pptx/antilag2降低觸控延遲.pptx')
    print("已生成：antilag2降低觸控延遲.pptx")

if __name__ == "__main__":
    create_pptx()
