# -*- coding: utf-8 -*-
"""
生成 Workload Prediction 一頁式投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 顏色定義
BG_COLOR = RGBColor(255, 249, 230)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_ORANGE = RGBColor(230, 126, 34)
ACCENT_GREEN = RGBColor(39, 174, 96)
ACCENT_RED = RGBColor(192, 57, 43)

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft JhengHei"

    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.16), height - Inches(0.32))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(1)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()

    # ===== 標題區 =====
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Workload Prediction - 降低功耗與提升 1% Low"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft JhengHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目標：從引擎底層取得每幀算力提示，實現精準控頻且不掉幀"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 左上：歷史問題 =====
    add_content_box(
        slide, Inches(0.25), Inches(0.85), Inches(4.25), Inches(2.3),
        "歷史問題回顧",
        [
            "1% Low 掉幀的根本原因：",
            "",
            "• CPU 在突發重載時來不及提頻",
            "• 為追求最佳功耗，給過高的頻率/算力",
            "• 導致功耗過高",
            "",
            "MTK 現行作法：",
            "• 人工針對遊戲調整頻率(算力)天花板",
            "• 問題：突發重載時常造成 1% low 掉幀",
            "",
            "需要：更精準的算力預測機制",
        ],
        title_color=ACCENT_RED
    )

    # ===== 中上：研究方法 =====
    add_content_box(
        slide, Inches(4.55), Inches(0.85), Inches(4.25), Inches(2.3),
        "研究方法",
        [
            "核心思路：",
            "修改 Unreal / Unity 引擎代碼，",
            "從引擎底層取得每幀的算力提示",
            "",
            "實作方式：",
            "• MAGT 製作 Unreal Demo App 模擬負載",
            "• App 能模擬鳴潮的負載特性",
            "• App 給出每幀算力需求 → 平台",
            "• 平台達到精準控頻且不掉幀",
            "",
            "實際遊戲：鳴潮 POC（已取得工作室同意）",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 右上：目標遊戲 =====
    add_content_box(
        slide, Inches(8.85), Inches(0.85), Inches(4.2), Inches(2.3),
        "目標遊戲與場景",
        [
            "重載遊戲驗證：",
            "• 鳴潮 / Demo App",
            "",
            "高刷遊戲驗證 (165 FPS)：",
            "• Demo App",
            "",
            "候選高刷 Unreal 遊戲：",
            "• 和平菁英",
            "• 三角洲行動",
            "• 暗區突圍",
            "• 無畏契約",
            "（擇一進行驗證）",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 中間左：POC 設計 =====
    add_content_box(
        slide, Inches(0.25), Inches(3.25), Inches(6.5), Inches(1.85),
        "POC 設計",
        [
            "測試場景：",
            "① 重載遊戲（鳴潮 / Demo App）",
            "② 重載遊戲高刷（165 FPS / Demo App）",
            "",
            "基準條件：",
            "• 拉到最高檔吹風扇 → 穩定 1% low = 59.9 FPS",
            "",
            "實驗條件（使用 Workload Prediction 演算法）：",
            "• 不吹風扇、板溫 < 38°C → 達到同樣 1% low 效果",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 中間右：技術架構 =====
    add_content_box(
        slide, Inches(6.85), Inches(3.25), Inches(6.2), Inches(1.85),
        "技術架構示意",
        [
            "引擎層（Unreal/Unity）",
            "    ↓ 每幀算力需求提示",
            "MAGT SDK",
            "    ↓ Workload Hint",
            "平台層（控頻策略）",
            "    ↓ 精準調頻",
            "硬體層（CPU/GPU）",
            "",
            "關鍵：引擎主動告知「下一幀需要多少算力」",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 底部左：預期結果 =====
    add_content_box(
        slide, Inches(0.25), Inches(5.2), Inches(6.5), Inches(2.05),
        "預期結果",
        [
            "1. 功耗下降",
            "   因減少不必要的急拉頻率",
            "",
            "2. 1% Low 提升至接近 Avg FPS",
            "   因可以及時急拉（精準預測）",
            "",
            "3. 功耗 ≤ 原演算法",
            "   在達到同樣效能的前提下",
            "   整體功耗不增加（理想是下降）",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 底部右：成功判定 =====
    add_content_box(
        slide, Inches(6.85), Inches(5.2), Inches(6.2), Inches(2.05),
        "成功判定總結",
        [
            "場景：重載遊戲 / 高刷遊戲",
            "",
            "指標：",
            "• 1% Low：達到 59.9 FPS（與吹風扇基準相當）",
            "• 溫度：板溫 < 38°C（無需吹風扇散熱）",
            "• 功耗：≤ 原演算法",
            "",
            "核心價值：",
            "透過引擎級算力預測，同時達成「省電」與「不掉幀」",
            "解決過去「功耗 vs 1% low」的兩難困境",
        ],
        title_color=ACCENT_BLUE
    )

    prs.save('D:/pptx/Workload預測降低功耗與1%Low.pptx')
    print("已生成：Workload預測降低功耗與1%Low.pptx")

if __name__ == "__main__":
    create_pptx()
