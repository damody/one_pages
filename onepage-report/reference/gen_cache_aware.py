# -*- coding: utf-8 -*-
"""
生成 Cache Aware Scheduling 一頁式投影片
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 顏色定義
BG_COLOR = RGBColor(255, 249, 230)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ACCENT_BLUE = RGBColor(70, 130, 180)
ACCENT_ORANGE = RGBColor(230, 126, 34)
ACCENT_GREEN = RGBColor(39, 174, 96)
ACCENT_PURPLE = RGBColor(142, 68, 173)

def set_cell_text(cell, text, font_size=10, bold=True, color=DARK_GRAY):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = "Microsoft JhengHei"

def add_content_box(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.03), width - Inches(0.16), Inches(0.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Microsoft JhengHei"

    content_box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.28), width - Inches(0.16), height - Inches(0.32))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(1)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()

    # ===== 標題區 =====
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cache Aware Scheduling - 解決掉幀問題"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft JhengHei"

    subtitle_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "靈感來源：Linux sched_cache (Intel/AMD) | 目標：降低跨 cluster 遷移，改善效能穩定度與功耗"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 左上：PC sched_cache 成功要素 =====
    add_content_box(
        slide, Inches(0.25), Inches(0.85), Inches(4.25), Inches(2.0),
        "PC sched_cache 成功要素",
        [
            "在 AMD EPYC Genoa 上驗證有效的前提：",
            "",
            "• 平台具備多個 cache domain (NUMA/CCX)",
            "• Thread 屬於長壽命、持續執行型",
            "• Thread 之間存在高度資料共享",
            "• 跨 cache domain 遷移有實質 cache 破壞成本",
            "",
            "效果：Host time -45%，Throughput +82%",
        ],
        title_color=ACCENT_BLUE
    )

    # ===== 中上：Dimensity Cache 架構 =====
    add_content_box(
        slide, Inches(4.55), Inches(0.85), Inches(4.25), Inches(2.0),
        "Dimensity 9500 Cache 架構",
        [
            "CPU 架構：1×Ultra + 3×Big + 4×Little",
            "",
            "Cache 層級：",
            "• 每個 cluster 內共用一顆 L2 cache",
            "• 不同 cluster 之間不共享 L2",
            "",
            "SoC 層級快取：",
            "• L3 cache：16MB / SLC：10MB（跨 cluster 共用）",
            "• LPDDR5X 主記憶體",
        ],
        title_color=ACCENT_PURPLE
    )

    # ===== 右上：為何可能有效 =====
    add_content_box(
        slide, Inches(8.85), Inches(0.85), Inches(4.2), Inches(2.0),
        "為何可能有類似效果",
        [
            "若以下條件在手機遊戲場景成立：",
            "",
            "• GameThread/RenderThread 長壽命、持續執行",
            "• Threads 間有穩定資料共享 (scene/render state)",
            "• 跨 cluster 遷移導致 L2 內容不可沿用",
            "• Working set 需重新 warm-up",
            "• LPDDR5X 延遲放大 cache miss 成本",
            "",
            "→ 降低跨 cluster migration 有機會改善穩定度",
        ],
        title_color=ACCENT_GREEN
    )

    # ===== 中間：對照表 =====
    table_title = slide.shapes.add_textbox(Inches(0.25), Inches(2.95), Inches(6), Inches(0.3))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "PC vs Mobile 結構相似性對照"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Microsoft JhengHei"

    table = slide.shapes.add_table(3, 3, Inches(0.25), Inches(3.2), Inches(6.5), Inches(1.0)).table
    table.columns[0].width = Inches(1.3)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.7)

    headers = ["平台", "Cache Domain 邊界", "特性"]
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, font_size=10, bold=True, color=WHITE)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = ACCENT_BLUE

    data = [
        ["AMD Genoa", "CCX/NUMA (L2不共享，每8核共享1個L3)", "跨 CCX 遷移成本高"],
        ["Dimensity", "Cluster (L2不共享，L3/SLC共用)", "跨 Cluster 遷移在 L2 層級即有成本"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            set_cell_text(table.cell(row_idx + 1, col_idx), cell_text, font_size=9, bold=True)
            if row_idx % 2 == 0:
                table.cell(row_idx + 1, col_idx).fill.solid()
                table.cell(row_idx + 1, col_idx).fill.fore_color.rgb = RGBColor(245, 245, 245)

    # 說明文字
    note_box = slide.shapes.add_textbox(Inches(0.25), Inches(4.25), Inches(6.5), Inches(0.4))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dimensity 非與 Genoa 等價，但在「跨 domain 遷移有實質成本」上具可比性，且手機 L2 邊界更早、容量更小、對 IPC 更敏感"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.font.name = "Microsoft JhengHei"

    # ===== 右中：MAGT 整合 =====
    add_content_box(
        slide, Inches(6.85), Inches(2.95), Inches(6.2), Inches(1.7),
        "MAGT 整合方案",
        [
            "功能開關：",
            "• 透過 MAGT 下 hint 開關此 Feature",
            "• 僅在有效益的遊戲場景啟用",
            "",
            "配合排程器：",
            "• 現行方案 (Loom) + LAVD (cache aware)",
            "• 讓 GameThread/RenderThread 盡量留在同 cluster",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 底部左：POC 設計 =====
    add_content_box(
        slide, Inches(0.25), Inches(4.7), Inches(6.5), Inches(2.55),
        "POC 設計",
        [
            "實驗條件：",
            "• A：現行方案 (FPSGO/Loom)",
            "• B：現行方案 (Loom) + LAVD (cache aware)",
            "",
            "遊戲場景：",
            "① 遊戲高刷（穩態）",
            "② 遊戲 + 抖音複合場景（干擾/搶佔）",
            "",
            "觀測指標：",
            "• 關鍵 thread 的 cluster 遷移次數",
            "• L2 cache miss rate 變化",
        ],
        title_color=ACCENT_ORANGE
    )

    # ===== 底部右：成功判定 =====
    add_content_box(
        slide, Inches(6.85), Inches(4.7), Inches(6.2), Inches(2.55),
        "成功判定準則",
        [
            "在 Avg FPS 相近（或 target FPS 相同）前提下：",
            "",
            "1. 1% Low 明顯提升",
            "   Frame Time 穩定度改善",
            "",
            "2. 關鍵 Thread 延遲下降",
            "   wakeup latency / runnable delay 顯著下降",
            "",
            "3. 功耗波動不惡化",
            "   理想是下降；至少不因救 delay 而急拉頻率",
            "",
            "合理推論：降低不必要跨 cluster migration 有機會帶來 PC 類似方向效益",
        ],
        title_color=ACCENT_BLUE
    )

    prs.save('D:/pptx/cache_aware排程解決掉幀問題.pptx')
    print("已生成：cache_aware排程解決掉幀問題.pptx")

if __name__ == "__main__":
    create_pptx()
