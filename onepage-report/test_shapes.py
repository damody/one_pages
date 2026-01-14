"""
測試 PPTX Shapes 繪圖功能
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# === 顏色定義 ===
COLOR_RED = RGBColor(244, 67, 54)
COLOR_GREEN = RGBColor(76, 175, 80)
COLOR_BLUE = RGBColor(33, 150, 243)
COLOR_GRAY_BG = RGBColor(245, 245, 245)
COLOR_TEXT = RGBColor(51, 51, 51)
COLOR_WHITE = RGBColor(255, 255, 255)


# === 前後對比圖 ===
def draw_before_after(slide, left, top, width, height, before_title, before_items, after_title, after_items):
    """繪製前後對比圖"""
    box_width = (width - 0.4) / 2

    # 左側區塊（改善前）
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(box_width), Inches(height)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = COLOR_GRAY_BG
    before_box.line.color.rgb = COLOR_RED
    before_box.line.width = Pt(2)

    # 左側標題
    before_title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = before_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = before_title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    # 左側內容
    before_content = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.45),
        Inches(box_width - 0.2), Inches(height - 0.55)
    )
    tf = before_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(before_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT

    # 右側區塊（改善後）
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + box_width + 0.4), Inches(top),
        Inches(box_width), Inches(height)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = COLOR_GRAY_BG
    after_box.line.color.rgb = COLOR_GREEN
    after_box.line.width = Pt(2)

    # 右側標題
    after_title_box = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.1),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = after_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = after_title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    # 右側內容
    after_content = slide.shapes.add_textbox(
        Inches(left + box_width + 0.5), Inches(top + 0.45),
        Inches(box_width - 0.2), Inches(height - 0.55)
    )
    tf = after_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(after_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT

    # 中間箭頭
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left + box_width + 0.1), Inches(top + height/2 - 0.15),
        Inches(0.2), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_BLUE
    arrow.line.fill.background()


# === 流程圖 ===
def draw_flow(slide, left, top, width, height, nodes):
    """繪製橫向流程圖"""
    node_count = len(nodes)
    gap = 0.15
    arrow_width = 0.12
    node_width = (width - gap * (node_count - 1)) / node_count

    for i, node in enumerate(nodes):
        x = left + i * (node_width + gap)

        if isinstance(node, dict):
            title = node.get("title", "")
            desc = node.get("desc", "")
        else:
            title = str(node)
            desc = ""

        # 節點矩形
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top),
            Inches(node_width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BLUE
        shape.line.fill.background()

        # 節點文字
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(8)
            p2.font.color.rgb = COLOR_WHITE
            p2.alignment = PP_ALIGN.CENTER

        # 箭頭
        if i < node_count - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + node_width + 0.02), Inches(top + height/2 - 0.08),
                Inches(arrow_width), Inches(0.16)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.fill.background()


# === 平台對比圖 ===
def draw_platform_compare(slide, left, top, width, height, platform_a, platform_b):
    """繪製上下平台對比圖"""
    box_height = (height - 0.3) / 2

    for i, platform in enumerate([platform_a, platform_b]):
        y = top if i == 0 else top + box_height + 0.3
        color = platform.get("color", COLOR_BLUE)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(y),
            Inches(width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_GRAY_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)

        title_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(y + 0.1),
            Inches(width - 0.2), Inches(0.25)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = platform.get("name", "")
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        content_box = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(y + 0.4),
            Inches(width - 0.2), Inches(box_height - 0.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(platform.get("items", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_TEXT


# === 時間軸圖 ===
def draw_timeline(slide, left, top, width, height, stages):
    """繪製時間軸圖"""
    line_y = top + height * 0.6
    stage_count = len(stages)
    stage_width = width / stage_count

    # 水平主軸線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(line_y),
        Inches(width), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(100, 100, 100)
    line.line.fill.background()

    for i, stage in enumerate(stages):
        x = left + i * stage_width
        color = stage.get("color", COLOR_BLUE)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.05), Inches(top),
            Inches(stage_width - 0.1), Inches(height * 0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stage.get("name", "")
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        if stage.get("time"):
            time_box = slide.shapes.add_textbox(
                Inches(x), Inches(line_y + 0.08),
                Inches(stage_width), Inches(0.25)
            )
            tf = time_box.text_frame
            p = tf.paragraphs[0]
            p.text = stage["time"]
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.CENTER


# === 主程式 ===
if __name__ == "__main__":
    # 建立簡報
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === 第 1 頁：前後對比圖 ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "測試 1：前後對比圖 (before_after)"
    p.font.size = Pt(24)
    p.font.bold = True

    draw_before_after(
        slide=slide1,
        left=0.5, top=1.2, width=12.0, height=2.5,
        before_title="改善前：Frame Queue 堆積",
        before_items=[
            "GPU 不知道遊戲目標幀率",
            "BufferQueue 堆積 2-3 幀",
            "延遲累積達 50-80ms"
        ],
        after_title="改善後：SDK 同步機制",
        after_items=[
            "SDK 通知目標幀率",
            "GPU 同步渲染節奏",
            "延遲降低 81%"
        ]
    )

    # === 第 2 頁：流程圖 ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "測試 2：流程圖 (flow)"
    p.font.size = Pt(24)
    p.font.bold = True

    draw_flow(
        slide=slide2,
        left=0.5, top=1.5, width=12.0, height=1.0,
        nodes=[
            {"title": "觸控輸入", "desc": "5ms"},
            {"title": "遊戲處理", "desc": "16ms"},
            {"title": "GPU 渲染", "desc": "8ms"},
            {"title": "顯示輸出", "desc": "8ms"}
        ]
    )

    # 第二個流程圖（純文字節點）
    draw_flow(
        slide=slide2,
        left=0.5, top=3.5, width=12.0, height=0.8,
        nodes=["需求分析", "設計", "開發", "測試", "部署"]
    )

    # === 第 3 頁：平台對比圖 ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "測試 3：平台對比圖 (platform_compare)"
    p.font.size = Pt(24)
    p.font.bold = True

    draw_platform_compare(
        slide=slide3,
        left=0.5, top=1.2, width=12.0, height=3.5,
        platform_a={
            "name": "PC 平台",
            "color": COLOR_BLUE,
            "items": [
                "Anti-Lag SDK 直接與 GPU Driver 通訊",
                "遊戲呼叫 SDK 設定目標幀率",
                "Driver 同步渲染節奏，消除 Frame Queue 堆積"
            ]
        },
        platform_b={
            "name": "手機平台",
            "color": COLOR_GREEN,
            "items": [
                "透過 SurfaceFlinger 管理 BufferQueue",
                "需要新的 SDK 介面通知目標幀率",
                "可借鏡 PC 方案的同步機制"
            ]
        }
    )

    # === 第 4 頁：時間軸圖 ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "測試 4：時間軸圖 (timeline)"
    p.font.size = Pt(24)
    p.font.bold = True

    draw_timeline(
        slide=slide4,
        left=0.5, top=1.5, width=12.0, height=1.5,
        stages=[
            {"name": "觸控", "time": "5ms", "color": COLOR_BLUE},
            {"name": "遊戲邏輯", "time": "16ms", "color": COLOR_BLUE},
            {"name": "GPU 渲染", "time": "8ms", "color": COLOR_GREEN},
            {"name": "合成", "time": "4ms", "color": COLOR_BLUE},
            {"name": "顯示", "time": "8ms", "color": COLOR_RED}
        ]
    )

    # 儲存
    output_path = "./test_shapes_output.pptx"
    prs.save(output_path)
    print(f"測試完成！已儲存至：{output_path}")
    print("請用 PowerPoint 開啟檢查 4 種圖表是否正確顯示。")
